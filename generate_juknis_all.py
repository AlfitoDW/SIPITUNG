"""
Merge semua juknis SIPITUNG menjadi satu file PPTX:
  - Ketua Tim Kerja: Rencana Aksi (Perencanaan)
  - Ketua Tim Kerja: Pengukuran Kinerja
  - Kabag Umum: Review & Persetujuan

Jalankan: python3 generate_juknis_all.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree
import copy

# ─── Palet warna ─────────────────────────────────────────────────────────────
BIRU_TUA   = RGBColor(0x1A, 0x37, 0x6C)
BIRU_MID   = RGBColor(0x1E, 0x4D, 0x9B)
BIRU_MUDA  = RGBColor(0x2E, 0x75, 0xB6)
HIJAU_TUA  = RGBColor(0x1B, 0x6B, 0x35)
HIJAU_MID  = RGBColor(0x27, 0xAE, 0x60)
HIJAU_MUDA = RGBColor(0xD5, 0xF5, 0xE3)
HIJAU_KAB  = RGBColor(0x0D, 0x47, 0x2E)
HIJAU_KAB2 = RGBColor(0x1A, 0x6B, 0x44)
HIJAU_KAB3 = RGBColor(0x2E, 0x8B, 0x57)
AKSEN      = RGBColor(0xF5, 0xA6, 0x23)
MERAH      = RGBColor(0xC0, 0x39, 0x2B)
AMBER_BG   = RGBColor(0xFF, 0xF3, 0xCD)
AMBER_BD   = RGBColor(0xF5, 0xA6, 0x23)
MERAH_BG   = RGBColor(0xFD, 0xED, 0xEC)
HIJAU_BG   = RGBColor(0xE8, 0xF5, 0xE9)
PUTIH      = RGBColor(0xFF, 0xFF, 0xFF)
ABU_MUDA   = RGBColor(0xF2, 0xF4, 0xF8)
ABU_TUA    = RGBColor(0x44, 0x44, 0x44)
PLACEHOLDER_BG = RGBColor(0xE8, 0xEE, 0xF8)

# ─── Setup ──────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def bs():
    """blank slide"""
    return prs.slides.add_slide(prs.slide_layouts[6])

# ─── Helpers ─────────────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.width = lw
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if line: s.line.color.rgb = line
    else:    s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h, sz=Pt(12), bold=False,
       color=ABU_TUA, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = sz; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return box

def hdr(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.1, fill=BIRU_TUA)
    rect(slide, 0, 1.1, 13.33, 0.06, fill=AKSEN)
    tb(slide, title, 0.35, 0.1, 12.5, 0.68, sz=Pt(22), bold=True, color=PUTIH)
    if subtitle:
        tb(slide, subtitle, 0.35, 0.7, 12.5, 0.38, sz=Pt(11),
           color=RGBColor(0xC8, 0xD8, 0xF0))

def ftr(slide, n, total, label="SIPITUNG — Petunjuk Teknis | LLDIKTI Wilayah III"):
    rect(slide, 0, 7.1, 13.33, 0.4, fill=BIRU_TUA)
    tb(slide, label, 0.35, 7.1, 10, 0.4, sz=Pt(9),
       color=RGBColor(0xC8, 0xD8, 0xF0))
    tb(slide, f"{n} / {total}", 12.5, 7.1, 0.8, 0.4, sz=Pt(9),
       color=PUTIH, align=PP_ALIGN.RIGHT)

def step_row(slide, num, title, desc, lx, ty, w=6.6, h=0.85, accent=BIRU_MID):
    rect(slide, lx, ty, w, h, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(0.75))
    rect(slide, lx, ty, 0.5, h, fill=accent)
    tb(slide, str(num), lx, ty+h/2-0.22, 0.5, 0.44,
       sz=Pt(17), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, title, lx+0.58, ty+0.05, w-0.65, 0.3, sz=Pt(10.5), bold=True, color=BIRU_TUA)
    tb(slide, desc,  lx+0.58, ty+0.36, w-0.65, h-0.42, sz=Pt(9.5), color=ABU_TUA, wrap=True)

def ph(slide, caption, lx, ty, w, h):
    rect(slide, lx, ty, w, h, fill=PLACEHOLDER_BG, line=BIRU_MUDA, lw=Pt(1.5))
    rect(slide, lx, ty, w, 0.38, fill=BIRU_MUDA)
    tb(slide, caption, lx+0.15, ty+0.02, w-0.2, 0.35, sz=Pt(10), bold=True, color=PUTIH)
    inner = caption.replace("Tampilan ","").replace("Tombol ","")
    tb(slide, f"[ SCREENSHOT {inner.upper()} ]\n\nTambahkan screenshot di sini",
       lx+0.3, ty+0.55, w-0.45, h-0.7,
       sz=Pt(10), color=BIRU_MUDA, align=PP_ALIGN.CENTER, italic=True)

def note(slide, text, lx, ty, w, h,
         bg=AMBER_BG, bd=AMBER_BD, fc=RGBColor(0x7D,0x60,0x08)):
    rect(slide, lx, ty, w, h, fill=bg, line=bd, lw=Pt(1))
    tb(slide, text, lx+0.15, ty+0.04, w-0.22, h-0.08, sz=Pt(9), color=fc, wrap=True)

def toc(slide, items, col1_color=BIRU_MID):
    for i, (num, ttl, dsc) in enumerate(items):
        col=i%2; row=i//2
        lx=0.4+col*6.5; ty=1.35+row*1.12
        rect(slide, lx, ty, 6.1, 1.0, fill=PUTIH,
             line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
        rect(slide, lx, ty, 0.7, 1.0, fill=col1_color)
        tb(slide, num, lx+0.05, ty+0.2, 0.6, 0.55,
           sz=Pt(16), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
        tb(slide, ttl, lx+0.78, ty+0.06, 5.2, 0.32, sz=Pt(11), bold=True, color=BIRU_TUA)
        tb(slide, dsc, lx+0.78, ty+0.38, 5.2, 0.55, sz=Pt(9), color=ABU_TUA, wrap=True)

def cover_pg(slide, nama_pengguna, modul_text, sub_text, accent=BIRU_MID):
    rect(slide, 0, 0, 13.33, 7.5, fill=BIRU_TUA)
    rect(slide, 0, 0, 4.5, 7.5, fill=RGBColor(0x0F,0x20,0x44))
    rect(slide, 4.5, 0, 0.08, 7.5, fill=AKSEN)
    rect(slide, 0.4, 0.4, 3.7, 1.6, fill=RGBColor(0x1A,0x37,0x6C), line=AKSEN, lw=Pt(1.5))
    tb(slide, "[LOGO LLDIKTI\nWILAYAH III]", 0.4, 0.4, 3.7, 1.6,
       sz=Pt(11), color=RGBColor(0xC8,0xD8,0xF0), align=PP_ALIGN.CENTER)
    tb(slide, "LLDIKTI WILAYAH III", 0.3, 2.3, 3.9, 0.55,
       sz=Pt(13), bold=True, color=AKSEN)
    tb(slide, "Kementerian Pendidikan Tinggi,\nSains dan Teknologi",
       0.3, 2.8, 3.9, 0.7, sz=Pt(10), color=RGBColor(0xC8,0xD8,0xF0))
    rect(slide, 0.3, 3.65, 3.9, 0.05, fill=AKSEN)
    tb(slide, "PETUNJUK TEKNIS", 0.3, 3.8, 3.9, 0.4, sz=Pt(11), bold=True,
       color=RGBColor(0xC8,0xD8,0xF0))
    tb(slide, "Penggunaan Aplikasi\nSIPITUNG", 0.3, 4.2, 3.9, 0.9,
       sz=Pt(14), bold=True, color=PUTIH)
    tb(slide, "Versi 1.0  |  April 2026", 0.3, 6.9, 3.9, 0.4,
       sz=Pt(9), color=RGBColor(0x88,0xA8,0xD0))
    tb(slide, "SIPITUNG", 4.9, 1.1, 8.0, 1.4, sz=Pt(52), bold=True, color=PUTIH)
    tb(slide, "Sistem Informasi Perencanaan & Kinerja", 4.9, 2.45, 8.0, 0.5,
       sz=Pt(16), color=AKSEN)
    rect(slide, 4.9, 3.05, 8.0, 0.05, fill=RGBColor(0x88,0xA8,0xD0))
    tb(slide, modul_text, 4.9, 3.2, 8.0, 0.5, sz=Pt(15), bold=True, color=PUTIH)
    tb(slide, sub_text,   4.9, 3.75, 8.0, 0.8, sz=Pt(13),
       color=RGBColor(0xC8,0xD8,0xF0))
    chips = [("Pengguna", nama_pengguna), ("Tahun","2026"), ("Status","Aktif")]
    for i, (lbl, val) in enumerate(chips):
        cx = 4.9 + i * 2.7
        rect(slide, cx, 4.7, 2.4, 0.65, fill=accent,
             line=RGBColor(0x88,0xA8,0xD0), lw=Pt(0.75))
        tb(slide, lbl, cx+0.1, 4.72, 2.2, 0.25, sz=Pt(8), bold=True, color=AKSEN)
        tb(slide, val, cx+0.1, 4.97, 2.2, 0.3,  sz=Pt(11), bold=True, color=PUTIH)

def section_divider(slide, section_num, title, subtitle, items, accent=BIRU_MID):
    """Slide pemisah antar bagian juknis."""
    rect(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0x0F,0x20,0x44))
    rect(slide, 0, 0, 0.5, 7.5, fill=accent)
    rect(slide, 0, 3.2, 13.33, 0.06, fill=AKSEN)
    tb(slide, f"BAGIAN {section_num}", 0.8, 1.2, 12.0, 0.65,
       sz=Pt(14), bold=True, color=AKSEN)
    tb(slide, title, 0.8, 1.85, 12.0, 1.2, sz=Pt(38), bold=True, color=PUTIH)
    tb(slide, subtitle, 0.8, 3.05, 12.0, 0.55, sz=Pt(15),
       color=RGBColor(0xC8,0xD8,0xF0))
    for i, item in enumerate(items):
        ix = 0.8 + i * 4.1
        rect(slide, ix, 3.9, 3.8, 1.0, fill=accent)
        tb(slide, item, ix+0.15, 3.95, 3.5, 0.9,
           sz=Pt(11), color=PUTIH, align=PP_ALIGN.CENTER, wrap=True)
    tb(slide, "SIPITUNG — Petunjuk Teknis | LLDIKTI Wilayah III",
       0.8, 6.9, 12.0, 0.45, sz=Pt(9), color=RGBColor(0x88,0xA8,0xD0))


# ════════════════════════════════════════════════════════════════════════════
# MASTER COVER (slide 1)
# ════════════════════════════════════════════════════════════════════════════
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0x0A,0x14,0x2C))
rect(slide, 0, 0, 13.33, 7.5, fill=BIRU_TUA)
rect(slide, 0, 0, 4.5, 7.5, fill=RGBColor(0x0A,0x14,0x2C))
rect(slide, 4.5, 0, 0.1, 7.5, fill=AKSEN)
rect(slide, 0.4, 0.4, 3.7, 1.7, fill=RGBColor(0x12,0x28,0x52), line=AKSEN, lw=Pt(1.5))
tb(slide, "[LOGO LLDIKTI\nWILAYAH III]", 0.4, 0.4, 3.7, 1.7,
   sz=Pt(11), color=RGBColor(0xC8,0xD8,0xF0), align=PP_ALIGN.CENTER)
tb(slide, "LLDIKTI WILAYAH III", 0.3, 2.35, 3.9, 0.55, sz=Pt(13), bold=True, color=AKSEN)
tb(slide, "Kementerian Pendidikan Tinggi,\nSains dan Teknologi",
   0.3, 2.9, 3.9, 0.7, sz=Pt(10), color=RGBColor(0xC8,0xD8,0xF0))
rect(slide, 0.3, 3.75, 3.9, 0.05, fill=AKSEN)
tb(slide, "BUKU PETUNJUK TEKNIS", 0.3, 3.9, 3.9, 0.42,
   sz=Pt(10), bold=True, color=RGBColor(0xC8,0xD8,0xF0))
tb(slide, "Penggunaan Aplikasi\nSIPITUNG", 0.3, 4.35, 3.9, 0.9,
   sz=Pt(14), bold=True, color=PUTIH)
tb(slide, "Versi 1.0  |  April 2026", 0.3, 6.9, 3.9, 0.4,
   sz=Pt(9), color=RGBColor(0x88,0xA8,0xD0))

tb(slide, "SIPITUNG", 5.0, 0.8, 8.0, 1.5, sz=Pt(58), bold=True, color=PUTIH)
tb(slide, "Sistem Informasi Perencanaan & Kinerja",
   5.0, 2.3, 8.0, 0.55, sz=Pt(17), color=AKSEN)
rect(slide, 5.0, 3.0, 8.0, 0.06, fill=RGBColor(0x88,0xA8,0xD0))
tb(slide, "Buku Petunjuk Teknis Lengkap", 5.0, 3.15, 8.0, 0.5,
   sz=Pt(16), bold=True, color=PUTIH)

sections = [
    ("I", "Ketua Tim Kerja\nRencana Aksi",     BIRU_MID),
    ("II","Ketua Tim Kerja\nPengukuran Kinerja",HIJAU_MID),
    ("III","Kabag Umum\nReview & Persetujuan",  HIJAU_KAB2),
]
for i, (num, lbl, col) in enumerate(sections):
    sx = 5.0 + i * 2.75
    rect(slide, sx, 3.8, 2.55, 1.65, fill=col)
    tb(slide, f"BAG. {num}", sx+0.15, 3.88, 2.3, 0.35,
       sz=Pt(9), bold=True, color=PUTIH)
    tb(slide, lbl, sx+0.15, 4.22, 2.3, 1.15,
       sz=Pt(12), bold=True, color=PUTIH, align=PP_ALIGN.LEFT, wrap=True)

tb(slide, "LLDIKTI Wilayah III  |  Jakarta  |  2026",
   5.0, 5.75, 8.0, 0.45, sz=Pt(10), color=RGBColor(0x88,0xA8,0xD0))
tb(slide, "1", 12.9, 7.2, 0.4, 0.3, sz=Pt(8), color=RGBColor(0x88,0xA8,0xD0),
   align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════════
# DAFTAR ISI MASTER (slide 2)
# ════════════════════════════════════════════════════════════════════════════
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "Daftar Isi Lengkap", "Tiga bagian juknis dalam satu dokumen")
ftr(slide, 2, "—")

sections_toc = [
    (BIRU_MID,   "BAGIAN I",  "Ketua Tim Kerja — Rencana Aksi",
     ["Login & Navigasi", "Mengisi Target TW I–IV", "Tambah Rencana Kegiatan",
      "Submit ke Kabag", "Revisi & Submit Ulang"]),
    (HIJAU_MID,  "BAGIAN II", "Ketua Tim Kerja — Pengukuran Kinerja",
     ["Navigasi & Periode", "Isi Realisasi IKU", "IKU Kolaborasi (Co-PIC)",
      "Submit Laporan", "Revisi & Status"]),
    (HIJAU_KAB2, "BAGIAN III","Kabag Umum — Review & Persetujuan",
     ["Login & Dashboard", "Review Rencana Aksi", "Setujui / Tolak RA",
      "Review Laporan", "Export PDF"]),
]
for i, (col, sec, ttl, pts) in enumerate(sections_toc):
    ty = 1.38 + i * 1.9
    rect(slide, 0.4, ty, 12.5, 1.82, fill=PUTIH,
         line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
    rect(slide, 0.4, ty, 12.5, 0.06, fill=col)
    rect(slide, 0.4, ty+0.06, 1.6, 1.76, fill=col)
    tb(slide, sec, 0.5, ty+0.18, 1.45, 0.42, sz=Pt(10), bold=True,
       color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, ttl, 0.6, ty+0.62, 1.3, 1.1, sz=Pt(10), bold=True,
       color=PUTIH, align=PP_ALIGN.CENTER, wrap=True)
    for j, pt in enumerate(pts):
        px = 2.15 + j * 2.2
        rect(slide, px, ty+0.15, 2.0, 1.52, fill=ABU_MUDA)
        tb(slide, f"• {pt}", px+0.12, ty+0.22, 1.8, 1.35,
           sz=Pt(9.5), color=BIRU_TUA, wrap=True)

# ════════════════════════════════════════════════════════════════════════════
# ══════════════════════ BAGIAN I: KETUA TIM — RENCANA AKSI ═══════════════════
# ════════════════════════════════════════════════════════════════════════════
T1_START = 3   # slide number offset
T1 = 12        # slides in this section

slide = bs()
section_divider(slide, "I", "Rencana Aksi", "Panduan Pengisian & Pengajuan — Ketua Tim Kerja",
                ["Login & Navigasi Menu", "Mengisi Target TW I–IV", "Menambah Rencana Kegiatan",
                 "Submit ke Kabag Umum"],
                accent=BIRU_MID)
ftr(slide, 3, "—")

# ── I-1: Cover ───────────────────────────────────────────────────────────────
slide = bs()
cover_pg(slide, "Ketua Tim Kerja",
         "Bagian I — Rencana Aksi",
         "Penyusunan, Pengisian Target & Submit ke Kabag Umum",
         accent=BIRU_MID)

# ── I-2: Pendahuluan & Tujuan ────────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "I-01  |  Pendahuluan & Tujuan", "Latar belakang dan maksud dokumen ini")
ftr(slide, 5, "—", "SIPITUNG — Bagian I: Ketua Tim Kerja (Rencana Aksi) | LLDIKTI Wilayah III")

tb(slide, "Latar Belakang", 0.4, 1.38, 6.1, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 0.4, 1.38, 0.06, 0.38, fill=AKSEN)
tb(slide,
   "SIPITUNG (Sistem Informasi Perencanaan & Kinerja) dikembangkan untuk mendukung "
   "proses perencanaan kerja di lingkungan LLDIKTI Wilayah III secara digital, "
   "transparan, dan terstruktur. Setiap Tim Kerja diwajibkan menyusun Rencana Aksi "
   "yang memuat target triwulan dan rencana kegiatan untuk setiap Indikator Kinerja "
   "Utama (IKU) yang menjadi tanggung jawab mereka.",
   0.5, 1.85, 6.0, 1.8, sz=Pt(10.5), color=ABU_TUA, wrap=True)

tb(slide, "Tujuan Juknis Bagian I", 0.4, 3.78, 6.1, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 0.4, 3.78, 0.06, 0.38, fill=AKSEN)
tujuan_items = [
    "Memberikan panduan penggunaan modul Penyusunan Rencana Aksi bagi Ketua Tim Kerja.",
    "Memastikan data yang diinput sesuai format dan ketentuan yang berlaku.",
    "Meminimalkan kesalahan dalam proses input dan pengajuan Rencana Aksi.",
    "Menjadi referensi saat terjadi kendala teknis di lapangan.",
]
for i, t in enumerate(tujuan_items):
    ty = 4.28 + i * 0.54
    rect(slide, 0.5, ty+0.07, 0.22, 0.22, fill=BIRU_MID)
    tb(slide, str(i+1), 0.5, ty+0.04, 0.22, 0.28,
       sz=Pt(9), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, t, 0.82, ty, 5.7, 0.5, sz=Pt(10.5), color=ABU_TUA, wrap=True)

rect(slide, 7.0, 1.28, 6.0, 5.85, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
rect(slide, 7.0, 1.28, 6.0, 0.08, fill=BIRU_MUDA)
tb(slide, "Informasi Dokumen", 7.2, 1.4, 5.6, 0.38, sz=Pt(12), bold=True, color=BIRU_TUA)
doc_info = [("Nama Sistem","SIPITUNG"), ("Versi","1.0"), ("Tanggal Terbit","April 2026"),
            ("Pengguna","Ketua Tim Kerja"), ("Modul","Perencanaan > Rencana Aksi > Penyusunan"),
            ("Instansi","LLDIKTI Wilayah III"), ("Kontak","Admin Sistem — ext. 100")]
for i, (k, v) in enumerate(doc_info):
    ty = 1.92 + i * 0.7
    rect(slide, 7.0, ty, 6.0, 0.7, fill=ABU_MUDA if i%2==0 else PUTIH)
    tb(slide, k, 7.15, ty+0.1, 2.4, 0.52, sz=Pt(10), bold=True, color=BIRU_TUA)
    tb(slide, v, 9.6, ty+0.1, 3.3, 0.52, sz=Pt(10), color=ABU_TUA)

# ── I-3: Ketentuan Umum ──────────────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "I-02  |  Ketentuan Umum & Ruang Lingkup",
    "Batasan pengguna dan aturan yang wajib diketahui")
ftr(slide, 6, "—", "SIPITUNG — Bagian I: Ketua Tim Kerja (Rencana Aksi) | LLDIKTI Wilayah III")

ku_items = [
    ("1","Akun Aktif","Pengguna harus memiliki akun dengan status aktif. Hubungi Admin jika tidak dapat login."),
    ("2","Kelengkapan Data","Semua IKU yang tampil wajib diisi target TW I–IV sebelum dokumen dapat disubmit."),
    ("3","Dokumen Tersubmit","Dokumen yang telah disubmit tidak dapat diubah, kecuali dikembalikan oleh Kabag Umum."),
    ("4","IKU Kolaborasi","Target pada IKU kolaborasi diisi bersama, namun masing-masing tim mengisi rencana kegiatannya."),
    ("5","Kerahasiaan Akun","Username dan password bersifat rahasia. Jangan berbagi akun dengan pihak lain."),
    ("6","Format Username","Username Ketua Tim Kerja mengikuti format ketua.xxx (contoh: ketua.pk, ketua.hmk)."),
]
for i, (num, ttl, dsc) in enumerate(ku_items):
    col=i%2; row=i//2
    lx=0.4+col*6.5; ty=1.38+row*1.92
    rect(slide, lx, ty, 6.1, 1.85, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
    rect(slide, lx, ty, 6.1, 0.08, fill=BIRU_MID)
    rect(slide, lx, ty+0.08, 0.55, 1.77, fill=RGBColor(0xDB,0xE8,0xF8))
    tb(slide, num, lx+0.1, ty+0.62, 0.38, 0.55, sz=Pt(22), bold=True,
       color=BIRU_TUA, align=PP_ALIGN.CENTER)
    tb(slide, ttl, lx+0.65, ty+0.12, 5.3, 0.32, sz=Pt(11), bold=True, color=BIRU_TUA)
    rect(slide, lx+0.65, ty+0.46, 5.35, 0.02, fill=RGBColor(0xCC,0xD9,0xEA))
    tb(slide, dsc, lx+0.65, ty+0.52, 5.35, 1.26, sz=Pt(9.5), color=ABU_TUA, wrap=True)

# ── I-4: Login & Navigasi ────────────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "I-03  |  Login & Navigasi Menu", "Masuk ke aplikasi dan membuka halaman Penyusunan Rencana Aksi")
ftr(slide, 7, "—", "SIPITUNG — Bagian I: Ketua Tim Kerja (Rencana Aksi) | LLDIKTI Wilayah III")

rect(slide, 0.4, 1.35, 12.5, 0.52, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
for i, bc in enumerate(["Dashboard","Perencanaan","Rencana Aksi","Penyusunan"]):
    bx = 0.5+i*3.1
    bg = BIRU_TUA if i==3 else (BIRU_MUDA if i>0 else RGBColor(0xE0,0xE8,0xF5))
    fc = PUTIH if i>=1 else BIRU_TUA
    rect(slide, bx, 1.38, 2.8, 0.46, fill=bg)
    tb(slide, ("▶  " if i>0 else "")+bc, bx+0.1, 1.41, 2.6, 0.4,
       sz=Pt(10), bold=(i==3), color=fc)

ph(slide, "Tampilan Login & Sidebar Menu", 0.4, 2.0, 4.5, 5.0)

tb(slide, "Langkah Login & Navigasi", 5.25, 2.02, 7.7, 0.38,
   sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.25, 2.02, 0.06, 0.38, fill=AKSEN)
login_steps = [
    ("1","Buka & Login",'Buka browser, akses URL SIPITUNG. Ketik username (ketua.xxx) dan password, klik "Masuk".'),
    ("2","Klik Perencanaan","Pada sidebar kiri, klik menu 'Perencanaan' untuk memperluas submenu."),
    ("3","Pilih Rencana Aksi","Dari submenu, klik 'Rencana Aksi'."),
    ("4","Pilih Penyusunan","Klik 'Penyusunan' untuk masuk ke halaman daftar IKU tim Anda."),
    ("5","Verifikasi","Pastikan judul halaman menampilkan nama Tim Kerja Anda dan daftar IKU tampil."),
]
for i, (num, ttl, dsc) in enumerate(login_steps):
    step_row(slide, num, ttl, dsc, 5.25, 2.55+i*0.9, w=7.7, h=0.83, accent=BIRU_MID)

# ── I-5: Mengisi Target TW ───────────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "I-04  |  Mengisi Target Triwulan (TW I–IV)",
    "Klik ikon pensil pada setiap IKU untuk membuka form isian")
ftr(slide, 8, "—", "SIPITUNG — Bagian I: Ketua Tim Kerja (Rencana Aksi) | LLDIKTI Wilayah III")

ph(slide, "Tampilan Form Isian Target", 0.4, 1.35, 5.3, 5.75)

tb(slide, "Cara Mengisi Target Triwulan", 6.05, 1.38, 7.0, 0.38,
   sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 6.05, 1.38, 0.06, 0.38, fill=AKSEN)
target_steps = [
    ("1","Temukan IKU di Daftar","Pada halaman Penyusunan, lihat daftar IKU mandiri dan kolaborasi tim Anda."),
    ("2","Klik Ikon Pensil (✏)","Klik ikon edit pada baris IKU yang ingin diisi. Form isian akan terbuka."),
    ("3","Isi Target TW I–IV","Isi nilai target:\n• TW I: Target akhir Maret\n• TW II: Target akhir Juni\n• TW III: Target akhir September\n• TW IV: Target akhir Desember"),
    ("4","Klik Simpan",'Klik "Simpan" untuk menyimpan data. Status IKU berubah menandai sudah diisi.'),
    ("5","Ulangi untuk Semua IKU","Lakukan langkah 2–4 untuk setiap IKU yang tampil di daftar."),
]
for i, (num, ttl, dsc) in enumerate(target_steps):
    step_row(slide, num, ttl, dsc, 6.05, 1.9+i*0.97, w=7.0, h=0.9, accent=BIRU_MID)

note(slide, "✓  Tips: Target triwulan bersifat kumulatif. Pastikan nilai TW IV sesuai dengan target tahunan yang telah ditetapkan.",
     6.05, 6.82, 7.0, 0.55, bg=HIJAU_BG, bd=HIJAU_MID, fc=HIJAU_TUA)

# ── I-6: Menambah Rencana Kegiatan ──────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "I-05  |  Menambahkan Rencana Kegiatan",
    "Klik ikon list pada setiap IKU untuk mengelola kegiatan")
ftr(slide, 9, "—", "SIPITUNG — Bagian I: Ketua Tim Kerja (Rencana Aksi) | LLDIKTI Wilayah III")

ph(slide, "Tampilan Form Rencana Kegiatan", 0.4, 1.35, 5.3, 5.75)

tb(slide, "Cara Menambahkan Rencana Kegiatan", 6.05, 1.38, 7.0, 0.38,
   sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 6.05, 1.38, 0.06, 0.38, fill=AKSEN)
kegiatan_steps = [
    ("1","Klik Ikon Kegiatan (☰)","Pada baris IKU, klik ikon list/kegiatan untuk membuka panel rencana kegiatan."),
    ("2","Klik 'Tambah Kegiatan'","Klik tombol Tambah Kegiatan untuk membuka form isian kegiatan baru."),
    ("3","Isi Nama Kegiatan","Tuliskan nama kegiatan secara jelas dan spesifik."),
    ("4","Pilih Triwulan","Pilih triwulan pelaksanaan kegiatan (TW I, II, III, atau IV)."),
    ("5","Simpan Kegiatan",'Klik "Simpan". Kegiatan muncul dalam daftar rencana kegiatan IKU tersebut.'),
    ("6","Ulangi Sesuai Kebutuhan","Tambahkan semua kegiatan yang direncanakan per IKU per triwulan."),
]
for i, (num, ttl, dsc) in enumerate(kegiatan_steps):
    step_row(slide, num, ttl, dsc, 6.05, 1.9+i*0.82, w=7.0, h=0.76, accent=BIRU_MID)

note(slide, "⚠  Perhatian: Minimal satu kegiatan harus ditambahkan per IKU sebelum dokumen dapat disubmit.",
     6.05, 6.82, 7.0, 0.55)

# ── I-7: Submit & Revisi ─────────────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "I-06  |  Submit & Revisi Rencana Aksi",
    "Mengajukan dokumen ke Kabag Umum dan prosedur jika dikembalikan")
ftr(slide, 10, "—", "SIPITUNG — Bagian I: Ketua Tim Kerja (Rencana Aksi) | LLDIKTI Wilayah III")

alur_submit = [("Isi Semua\nIKU",BIRU_MUDA,"Tim Kerja"), ("Submit ke\nKabag",BIRU_MID,"Tim Kerja"),
               ("Review\nKabag",BIRU_TUA,"Kabag Umum"), ("Disetujui /\nDikembalikan",HIJAU_MID,"Kabag Umum")]
rect(slide, 0.4, 1.35, 12.5, 1.4, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
for i, (lbl, col, actor) in enumerate(alur_submit):
    ax = 0.6+i*3.1
    rect(slide, ax, 1.42, 2.7, 1.26, fill=col)
    tb(slide, lbl, ax+0.1, 1.47, 2.5, 0.6, sz=Pt(10), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, actor, ax+0.1, 2.04, 2.5, 0.55, sz=Pt(8.5), color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)
    if i<len(alur_submit)-1:
        tb(slide, "▶", ax+2.73, 1.84, 0.35, 0.5, sz=Pt(14), color=BIRU_MUDA, align=PP_ALIGN.CENTER)

rect(slide, 0.4, 2.9, 5.9, 4.15, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
rect(slide, 0.4, 2.9, 5.9, 0.42, fill=BIRU_TUA)
tb(slide, "Cara Submit Dokumen", 0.55, 2.92, 5.65, 0.38, sz=Pt(12), bold=True, color=PUTIH)
submit_ss = [
    ("1","Periksa Kelengkapan","Semua IKU sudah diisi target TW I–IV dan minimal satu rencana kegiatan."),
    ("2","Klik Tombol Submit",'Klik "Submit ke Kabag Umum" di halaman Penyusunan.'),
    ("3","Konfirmasi","Baca dialog konfirmasi dengan seksama, lalu klik Konfirmasi."),
    ("4","Status Berubah","Status dokumen menjadi 'Menunggu Persetujuan'."),
]
for i, (num, ttl, dsc) in enumerate(submit_ss):
    step_row(slide, num, ttl, dsc, 0.4, 3.42+i*0.83, w=5.9, h=0.76, accent=BIRU_TUA)

note(slide, "⛔  Dokumen yang telah disubmit TIDAK DAPAT diedit. Pastikan data sudah benar.",
     0.4, 6.88, 5.9, 0.55, bg=MERAH_BG, bd=MERAH, fc=MERAH)

rect(slide, 6.7, 2.9, 6.25, 4.15, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
rect(slide, 6.7, 2.9, 6.25, 0.42, fill=MERAH)
tb(slide, "Jika Dokumen Dikembalikan", 6.85, 2.92, 6.0, 0.38, sz=Pt(12), bold=True, color=PUTIH)
revisi_ss = [
    ("1","Baca Catatan","Cek banner merah di halaman Penyusunan. Klik 'Lihat Catatan' untuk membaca keterangan Kabag."),
    ("2","Perbaiki Data","Klik ikon pensil (✏) pada IKU yang perlu diperbaiki, ubah sesuai catatan."),
    ("3","Simpan","Simpan perubahan pada setiap IKU yang diperbaiki."),
    ("4","Submit Ulang",'Klik kembali "Submit ke Kabag Umum" setelah semua perbaikan selesai.'),
]
for i, (num, ttl, dsc) in enumerate(revisi_ss):
    step_row(slide, num, ttl, dsc, 6.7, 3.42+i*0.83, w=6.25, h=0.76, accent=MERAH)

note(slide, "⚠  Tidak ada batasan jumlah submit ulang. Pastikan semua catatan revisi sudah ditangani.",
     6.7, 6.88, 6.25, 0.55)

# ── I-8: FAQ ─────────────────────────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "I-07  |  FAQ — Rencana Aksi", "Pertanyaan yang Sering Diajukan")
ftr(slide, 11, "—", "SIPITUNG — Bagian I: Ketua Tim Kerja (Rencana Aksi) | LLDIKTI Wilayah III")

faqs_ra = [
    ("Saya tidak bisa login. Apa yang harus dilakukan?",
     "Periksa username (ketua.xxx) dan password. Pastikan Caps Lock nonaktif. Hubungi Admin untuk reset password."),
    ("Mengapa tombol Submit tidak aktif?",
     "Tombol Submit hanya aktif jika SEMUA IKU sudah diisi target TW I–IV dan minimal satu rencana kegiatan."),
    ("Bisakah saya mengedit data setelah submit?",
     "Tidak. Pengeditan hanya bisa dilakukan jika Kabag Umum mengembalikan dokumen dengan catatan revisi."),
    ("Apa itu IKU Kolaborasi? Siapa yang mengisinya?",
     "IKU Kolaborasi dikerjakan bersama beberapa Tim Kerja. Setiap Tim Kerja mengisi rencana kegiatannya masing-masing."),
    ("Berapa lama proses persetujuan oleh Kabag Umum?",
     "Tergantung jadwal Kabag Umum. Pantau status dokumen secara berkala di halaman Penyusunan."),
    ("Data yang saya input tidak tersimpan. Apa yang terjadi?",
     "Pastikan Anda selalu menekan tombol Simpan setelah mengisi data. Data tidak tersimpan otomatis."),
]
for i, (q, a) in enumerate(faqs_ra):
    col=i%2; row=i//2
    lx=0.4+col*6.5; ty=1.38+row*1.92
    rect(slide, lx, ty, 6.1, 1.85, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
    rect(slide, lx, ty, 6.1, 0.08, fill=BIRU_MID)
    rect(slide, lx, ty+0.08, 0.52, 1.77, fill=RGBColor(0xDB,0xE8,0xF8))
    tb(slide, "Q", lx+0.1, ty+0.66, 0.34, 0.44, sz=Pt(18), bold=True,
       color=BIRU_TUA, align=PP_ALIGN.CENTER)
    tb(slide, q, lx+0.62, ty+0.1, 5.35, 0.48, sz=Pt(10), bold=True, color=BIRU_TUA, wrap=True)
    rect(slide, lx+0.62, ty+0.6, 5.35, 0.02, fill=RGBColor(0xCC,0xD9,0xEA))
    tb(slide, a, lx+0.62, ty+0.68, 5.35, 1.12, sz=Pt(9.5), color=ABU_TUA, wrap=True)

# ════════════════════════════════════════════════════════════════════════════
# ══════════════════ BAGIAN II: KETUA TIM — PENGUKURAN KINERJA ════════════════
# ════════════════════════════════════════════════════════════════════════════
slide = bs()
section_divider(slide, "II", "Pengukuran Kinerja",
                "Panduan Pengisian Realisasi & Pelaporan — Ketua Tim Kerja",
                ["Navigasi & Periode Aktif", "Mengisi Data Realisasi IKU",
                 "IKU Kolaborasi Co-PIC", "Submit & Revisi Laporan"],
                accent=HIJAU_MID)
ftr(slide, 12, "—")

# ── II-1: Cover ──────────────────────────────────────────────────────────────
slide = bs()
cover_pg(slide, "Ketua Tim Kerja",
         "Bagian II — Pengukuran Kinerja",
         "Pengisian Realisasi IKU & Submit Laporan Triwulan",
         accent=HIJAU_TUA)

# ── II-2: Pendahuluan ────────────────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "II-01  |  Pendahuluan & Siklus Triwulan", "Konteks modul Pengukuran Kinerja")
ftr(slide, 14, "—", "SIPITUNG — Bagian II: Ketua Tim Kerja (Pengukuran) | LLDIKTI Wilayah III")

tb(slide, "Tentang Modul Pengukuran Kinerja", 0.4, 1.38, 6.0, 0.38,
   sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 0.4, 1.38, 0.06, 0.38, fill=HIJAU_MID)
tb(slide,
   "Modul Pengukuran Kinerja memungkinkan setiap Ketua Tim Kerja mengisi data realisasi "
   "capaian IKU per triwulan. Data ini menjadi dasar evaluasi kinerja dan pelaporan kepada "
   "Kabag Umum. Pengisian dilakukan setelah Admin Sistem mengaktifkan periode triwulan.",
   0.5, 1.85, 5.9, 1.3, sz=Pt(10.5), color=ABU_TUA, wrap=True)

tb(slide, "Tujuan", 0.4, 3.3, 6.0, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 0.4, 3.3, 0.06, 0.38, fill=HIJAU_MID)
for i, g in enumerate([
    "Memandu Ketua Tim Kerja dalam mengisi data realisasi IKU per triwulan.",
    "Memastikan data capaian kinerja tercatat secara akurat dan tepat waktu.",
    "Menjelaskan alur koordinasi untuk IKU yang dikerjakan bersama (Co-PIC).",
    "Menjadi panduan dalam proses submit laporan hingga mendapat persetujuan.",
]):
    ty = 3.82+i*0.54
    rect(slide, 0.5, ty+0.07, 0.24, 0.24, fill=HIJAU_MID)
    tb(slide, str(i+1), 0.5, ty+0.04, 0.24, 0.28, sz=Pt(9), bold=True,
       color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, g, 0.82, ty, 5.6, 0.5, sz=Pt(10.5), color=ABU_TUA, wrap=True)

rect(slide, 7.0, 1.35, 6.0, 5.7, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
rect(slide, 7.0, 1.35, 6.0, 0.45, fill=HIJAU_TUA)
tb(slide, "Siklus Pengukuran Kinerja Per Tahun", 7.15, 1.37, 5.7, 0.4,
   sz=Pt(11), bold=True, color=PUTIH)
for i, (tw, bulan, ket) in enumerate([
    ("TW I",  "Januari – Maret",    "Laporan capaian Q1"),
    ("TW II", "April – Juni",       "Laporan capaian Q2"),
    ("TW III","Juli – September",   "Laporan capaian Q3"),
    ("TW IV", "Oktober – Desember", "Laporan capaian Q4 + evaluasi akhir tahun"),
]):
    ty = 1.95+i*1.22
    bg = HIJAU_MUDA if i%2==0 else PUTIH
    rect(slide, 7.0, ty, 6.0, 1.18, fill=bg)
    rect(slide, 7.0, ty, 1.0, 1.18, fill=HIJAU_MID if i%2==0 else HIJAU_TUA)
    tb(slide, tw, 7.0, ty+0.35, 1.0, 0.5, sz=Pt(13), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, bulan, 8.1, ty+0.1, 4.7, 0.35, sz=Pt(11), bold=True, color=BIRU_TUA)
    tb(slide, ket, 8.1, ty+0.48, 4.7, 0.6, sz=Pt(9.5), color=ABU_TUA, wrap=True)

# ── II-3: Navigasi & Ketentuan ───────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "II-02  |  Navigasi & Ketentuan Pengukuran",
    "Cara mengakses menu dan aturan pengisian realisasi")
ftr(slide, 15, "—", "SIPITUNG — Bagian II: Ketua Tim Kerja (Pengukuran) | LLDIKTI Wilayah III")

ph(slide, "Tampilan Menu Pengukuran Kinerja", 0.4, 1.35, 4.5, 5.75)

tb(slide, "Langkah Navigasi", 5.25, 1.38, 7.7, 0.38,
   sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.25, 1.38, 0.06, 0.38, fill=AKSEN)
for i, (num, ttl, dsc) in enumerate([
    ("1","Login","Masukkan username (ketua.xxx) dan password, lalu klik Masuk."),
    ("2","Klik Pengukuran Kinerja","Di sidebar kiri, klik menu 'Pengukuran Kinerja' (ikon grafik batang)."),
    ("3","Pilih Periode Aktif","Pilih periode triwulan aktif dari dropdown."),
    ("4","Verifikasi IKU","Pastikan daftar IKU tampil. Jika kosong, periode mungkin belum aktif."),
]):
    step_row(slide, num, ttl, dsc, 5.25, 1.9+i*0.88, w=7.7, h=0.82, accent=HIJAU_MID)

tb(slide, "Ketentuan Penting", 5.25, 5.5, 7.7, 0.38, sz=Pt(12), bold=True, color=BIRU_TUA)
rect(slide, 5.25, 5.5, 0.06, 0.38, fill=AKSEN)
for i, (k, v) in enumerate([
    ("Periode Aktif","Pengisian hanya bisa dilakukan saat Admin mengaktifkan periode triwulan."),
    ("IKU yang tampil","Hanya IKU di mana Tim Kerja Anda tercatat sebagai PIC yang akan tampil."),
    ("IKU Kolaborasi","Semua PIC dapat mengisi. Data terakhir yang disimpan menjadi nilai final (last-save-wins)."),
]):
    ty = 5.98+i*0.42
    rect(slide, 5.25, ty, 7.7, 0.4, fill=ABU_MUDA if i%2==0 else PUTIH)
    tb(slide, k, 5.38, ty+0.05, 1.8, 0.3, sz=Pt(9.5), bold=True, color=BIRU_TUA)
    tb(slide, v, 7.2, ty+0.05, 5.65, 0.3, sz=Pt(9.5), color=ABU_TUA)

# ── II-4: Mengisi Realisasi ───────────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "II-03  |  Mengisi Data Realisasi & IKU Kolaborasi",
    "Panduan pengisian nilai realisasi dan koordinasi Co-PIC")
ftr(slide, 16, "—", "SIPITUNG — Bagian II: Ketua Tim Kerja (Pengukuran) | LLDIKTI Wilayah III")

ph(slide, "Tampilan Form Isian Realisasi", 0.4, 1.35, 5.2, 5.75)

tb(slide, "Cara Mengisi Realisasi IKU", 5.95, 1.38, 7.1, 0.38,
   sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.95, 1.38, 0.06, 0.38, fill=AKSEN)
for i, (num, ttl, dsc) in enumerate([
    ("1","Klik 'Isi' atau 'Edit'","Pada baris IKU, klik tombol 'Isi' (belum ada data) atau 'Edit' (data sudah ada)."),
    ("2","Isi Nilai Realisasi (Wajib)","Masukkan nilai capaian aktual sesuai satuan IKU."),
    ("3","Isi Progress Kegiatan","Deskripsi singkat progres kegiatan yang telah dilaksanakan."),
    ("4","Isi Kendala & Strategi","Tuliskan kendala (jika ada) dan strategi tindak lanjut."),
    ("5","Klik Simpan",'Klik "Simpan". Status IKU berubah menjadi ✓ Terisi.'),
]):
    step_row(slide, num, ttl, dsc, 5.95, 1.9+i*0.85, w=7.1, h=0.78, accent=HIJAU_MID)

# IKU Kolaborasi box
rect(slide, 5.95, 6.2, 7.1, 1.3, fill=RGBColor(0xE3,0xF2,0xFD),
     line=BIRU_MUDA, lw=Pt(1))
rect(slide, 5.95, 6.2, 7.1, 0.38, fill=BIRU_MID)
tb(slide, "IKU Kolaborasi (Co-PIC)", 6.1, 6.22, 6.8, 0.34, sz=Pt(11), bold=True, color=PUTIH)
tb(slide, "• Ikon 👥 pada IKU menandakan IKU kolaborasi  • Semua PIC bisa isi/edit kapan saja\n"
         "• Gunakan kolom Catatan untuk koordinasi  • Hanya satu tim yang perlu submit laporan",
   6.1, 6.63, 6.8, 0.8, sz=Pt(9.5), color=BIRU_TUA, wrap=True)

# ── II-5: Submit & Revisi Laporan ─────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "II-04  |  Submit & Revisi Laporan Pengukuran",
    "Mengirim laporan ke Kabag Umum dan prosedur jika dikembalikan")
ftr(slide, 17, "—", "SIPITUNG — Bagian II: Ketua Tim Kerja (Pengukuran) | LLDIKTI Wilayah III")

for i, (lbl, col, actor) in enumerate([
    ("Isi Semua\nRealisasi",HIJAU_MID,"Tim Kerja"), ("Submit\nLaporan",BIRU_MID,"Tim Kerja"),
    ("Review\nKabag",BIRU_TUA,"Kabag Umum"), ("Disetujui /\nDikembalikan",HIJAU_TUA,"Kabag Umum")
]):
    ax = 0.6+i*3.1
    rect(slide, 0.4, 1.35, 12.5, 1.35, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
    rect(slide, ax, 1.42, 2.7, 1.2, fill=col)
    tb(slide, lbl, ax+0.1, 1.47, 2.5, 0.55, sz=Pt(10), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, actor, ax+0.1, 1.99, 2.5, 0.5, sz=Pt(8.5), color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)
    if i<3: tb(slide, "▶", ax+2.73, 1.78, 0.35, 0.5, sz=Pt(14), color=BIRU_MUDA, align=PP_ALIGN.CENTER)

rect(slide, 0.4, 2.85, 5.9, 4.15, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
rect(slide, 0.4, 2.85, 5.9, 0.42, fill=HIJAU_TUA)
tb(slide, "Cara Submit Laporan", 0.55, 2.87, 5.65, 0.38, sz=Pt(12), bold=True, color=PUTIH)
for i, (num, ttl, dsc) in enumerate([
    ("1","Periksa Progress Bar","Semua langkah: Periode ✓ | IKU tersedia ✓ | Realisasi diisi ✓."),
    ("2","Klik Submit Laporan",'Klik "Submit Laporan ke Kabag Umum".'),
    ("3","Konfirmasi","Klik Konfirmasi di dialog yang muncul."),
    ("4","Pantau Status","Status berubah menjadi 'Menunggu persetujuan Kabag Umum'."),
]):
    step_row(slide, num, ttl, dsc, 0.4, 3.37+i*0.83, w=5.9, h=0.76, accent=HIJAU_TUA)

note(slide, "⛔  Laporan yang telah disubmit TIDAK DAPAT diedit. Pastikan semua data sudah benar.",
     0.4, 6.83, 5.9, 0.55, bg=MERAH_BG, bd=MERAH, fc=MERAH)

rect(slide, 6.7, 2.85, 6.25, 4.15, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
rect(slide, 6.7, 2.85, 6.25, 0.42, fill=MERAH)
tb(slide, "Jika Laporan Dikembalikan", 6.85, 2.87, 6.0, 0.38, sz=Pt(12), bold=True, color=PUTIH)
for i, (num, ttl, dsc) in enumerate([
    ("1","Baca Catatan Revisi","Banner merah muncul di halaman Pengukuran. Baca rekomendasi Kabag."),
    ("2","Perbaiki Data","Klik 'Edit' pada IKU yang perlu diperbaiki. Ubah nilai sesuai catatan."),
    ("3","Simpan","Simpan perubahan pada setiap IKU yang diperbaiki."),
    ("4","Submit Ulang","Klik kembali tombol Submit Laporan. Catatan revisi sebelumnya reset otomatis."),
]):
    step_row(slide, num, ttl, dsc, 6.7, 3.37+i*0.83, w=6.25, h=0.76, accent=MERAH)

note(slide, "⚠  Status laporan: Belum Submit (kuning) → Menunggu Review (biru) → Disetujui (hijau) / Dikembalikan (merah).",
     6.7, 6.83, 6.25, 0.55)

# ── II-6: FAQ Pengukuran ─────────────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "II-05  |  FAQ — Pengukuran Kinerja", "Pertanyaan yang Sering Diajukan")
ftr(slide, 18, "—", "SIPITUNG — Bagian II: Ketua Tim Kerja (Pengukuran) | LLDIKTI Wilayah III")

for i, (q, a) in enumerate([
    ("IKU saya tidak tampil di halaman Pengukuran.",
     "Kemungkinan: periode belum diaktifkan oleh Admin, atau tidak ada IKU yang di-assign ke Tim Anda. Hubungi Admin."),
    ("Saya sudah mengisi tapi data hilang setelah refresh.",
     "Pastikan Anda sudah klik tombol 'Simpan'. Data tidak tersimpan otomatis."),
    ("Tombol Submit Laporan tidak aktif.",
     "Pastikan semua IKU sudah terisi (status ✓). Cek progress bar apakah ada step yang belum terpenuhi."),
    ("Tim co-PIC sudah submit lebih dulu. Apakah saya tetap harus submit?",
     "Tidak wajib. Namun Tim Anda tetap dapat submit secara mandiri jika diperlukan."),
    ("Nilai realisasi saya berbeda dengan milik co-PIC.",
     "Koordinasikan. Nilai terakhir yang disimpan menjadi nilai final (last-save-wins). Gunakan kolom Catatan."),
    ("Kabag mengembalikan laporan. Apakah data realisasi terhapus?",
     "Tidak. Data realisasi tetap tersimpan. Cukup perbaiki sesuai catatan dan submit ulang."),
]):
    col=i%2; row=i//2
    lx=0.4+col*6.5; ty=1.38+row*1.92
    rect(slide, lx, ty, 6.1, 1.85, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
    rect(slide, lx, ty, 6.1, 0.08, fill=HIJAU_MID)
    rect(slide, lx, ty+0.08, 0.52, 1.77, fill=HIJAU_MUDA)
    tb(slide, "Q", lx+0.1, ty+0.66, 0.34, 0.44, sz=Pt(18), bold=True,
       color=HIJAU_TUA, align=PP_ALIGN.CENTER)
    tb(slide, q, lx+0.62, ty+0.1, 5.35, 0.48, sz=Pt(10), bold=True, color=BIRU_TUA, wrap=True)
    rect(slide, lx+0.62, ty+0.6, 5.35, 0.02, fill=RGBColor(0xCC,0xD9,0xEA))
    tb(slide, a, lx+0.62, ty+0.68, 5.35, 1.12, sz=Pt(9.5), color=ABU_TUA, wrap=True)

# ════════════════════════════════════════════════════════════════════════════
# ══════════════════════ BAGIAN III: KABAG UMUM ════════════════════════════════
# ════════════════════════════════════════════════════════════════════════════
slide = bs()
section_divider(slide, "III", "Kabag Umum",
                "Panduan Review, Persetujuan & Export Laporan",
                ["Review Rencana Aksi", "Setujui / Tolak RA",
                 "Review Laporan Pengukuran", "Export PDF per TW"],
                accent=HIJAU_KAB2)
ftr(slide, 19, "—")

# ── III-1: Cover ─────────────────────────────────────────────────────────────
slide = bs()
cover_pg(slide, "Kabag Umum (Pimpinan)",
         "Bagian III — Review & Persetujuan",
         "Rencana Aksi · Laporan Pengukuran · Export PDF",
         accent=HIJAU_KAB)

# ── III-2: Peran & Ketentuan ─────────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "III-01  |  Peran & Ketentuan Kabag Umum",
    "Posisi dalam alur persetujuan dan aturan kewenangan")
ftr(slide, 21, "—", "SIPITUNG — Bagian III: Kabag Umum | LLDIKTI Wilayah III")

rect(slide, 0.4, 1.35, 12.5, 1.42, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
for i, (node, col, actor) in enumerate([
    ("Tim Kerja\n(Ketua)",BIRU_MID,"Menyusun RA &\nMengisi Realisasi"),
    ("Kabag Umum\n★ ANDA",HIJAU_KAB,"Review, Setujui/\nTolak Dokumen"),
    ("PPK",BIRU_TUA,"Persetujuan Final"),
    ("Terkunci",HIJAU_MID,"Dokumen Arsip"),
]):
    nx = 0.8+i*3.1
    rect(slide, nx, 1.42, 2.7, 1.28, fill=col)
    tb(slide, node, nx+0.1, 1.47, 2.5, 0.55, sz=Pt(11), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, actor, nx+0.1, 2.02, 2.5, 0.62, sz=Pt(9), color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)
    if i<3: tb(slide, "▶", nx+2.73, 1.82, 0.35, 0.5, sz=Pt(14), color=BIRU_MUDA, align=PP_ALIGN.CENTER)

rect(slide, 0.4, 2.92, 5.9, 4.12, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
rect(slide, 0.4, 2.92, 5.9, 0.42, fill=HIJAU_KAB2)
tb(slide, "Kewenangan Kabag Umum", 0.55, 2.94, 5.65, 0.38, sz=Pt(12), bold=True, color=PUTIH)
for i, k in enumerate([
    "Mereview dan menyetujui Rencana Aksi (RA) yang disubmit Tim Kerja",
    "Menolak RA disertai catatan revisi untuk diperbaiki Tim Kerja",
    "Mereview dan menyetujui Laporan Pengukuran Kinerja per triwulan",
    "Menolak Laporan disertai rekomendasi perbaikan",
    "Mengekspor laporan pengukuran dalam format PDF per triwulan",
]):
    ty = 3.44+i*0.5
    rect(slide, 0.4, ty, 0.42, 0.48, fill=HIJAU_KAB3)
    tb(slide, "✓", 0.4, ty+0.05, 0.42, 0.38, sz=Pt(13), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, k, 0.9, ty+0.06, 5.3, 0.4, sz=Pt(10), color=ABU_TUA, wrap=True)

rect(slide, 6.7, 2.92, 6.25, 4.12, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
rect(slide, 6.7, 2.92, 6.25, 0.42, fill=BIRU_TUA)
tb(slide, "Ketentuan Umum Kabag Umum", 6.85, 2.94, 6.0, 0.38, sz=Pt(12), bold=True, color=PUTIH)
for i, (k, v) in enumerate([
    ("Status Dokumen","Tombol Setujui/Tolak hanya aktif untuk dokumen berstatus 'submitted'."),
    ("Catatan Wajib","Setiap penolakan WAJIB disertai catatan/rekomendasi yang jelas."),
    ("View-Only","Halaman menampilkan SEMUA dokumen semua status. Tombol hanya untuk 'submitted'."),
    ("Persetujuan PPK","Setelah Kabag menyetujui, dokumen diteruskan ke PPK. Tidak bisa dibatalkan."),
    ("Export","Hanya Kabag Umum yang dapat mengekspor laporan ke PDF."),
]):
    ty = 3.44+i*0.72
    rect(slide, 6.7, ty, 6.25, 0.7, fill=ABU_MUDA if i%2==0 else PUTIH)
    tb(slide, k, 6.85, ty+0.07, 1.8, 0.3, sz=Pt(10), bold=True, color=BIRU_TUA)
    tb(slide, v, 8.65, ty+0.07, 4.2, 0.58, sz=Pt(9.5), color=ABU_TUA, wrap=True)

# ── III-3: Login & Review RA ─────────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "III-02  |  Login & Review Rencana Aksi",
    "Cara masuk dan mereview RA dari seluruh Tim Kerja")
ftr(slide, 22, "—", "SIPITUNG — Bagian III: Kabag Umum | LLDIKTI Wilayah III")

ph(slide, "Tampilan Halaman Review RA", 0.4, 1.35, 4.8, 5.75)

tb(slide, "Login & Navigasi Review RA", 5.55, 1.38, 7.4, 0.38,
   sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.55, 1.38, 0.06, 0.38, fill=AKSEN)
for i, (num, ttl, dsc) in enumerate([
    ("1","Login","Masukkan username kabag.umum dan password, lalu klik Masuk."),
    ("2","Buka Dashboard","Dashboard menampilkan ringkasan: Menunggu Review / Disetujui / Ditolak."),
    ("3","Menu Perencanaan → RA","Klik 'Perencanaan' di sidebar → 'Rencana Aksi'. Semua RA dari semua Tim Kerja tampil."),
    ("4","Identifikasi yang Perlu Direview","Cari baris dengan badge biru 'Menunggu Review'. Ini yang perlu ditindaklanjuti."),
    ("5","Review Kelengkapan","Periksa: target TW I–IV tidak kosong, nilai wajar, IKU kolaborasi ter-assign benar."),
    ("6","Ambil Keputusan","Jika data lengkap dan valid → Setujui. Jika perlu perbaikan → Tolak + catatan."),
]):
    step_row(slide, num, ttl, dsc, 5.55, 1.9+i*0.83, w=7.4, h=0.76, accent=HIJAU_KAB2)

note(slide, "ℹ  Tabel menampilkan SEMUA Tim Kerja & SEMUA status. Kolom Aksi (Setujui/Tolak) hanya aktif untuk status 'submitted'.",
     5.55, 6.95, 7.4, 0.48, bg=RGBColor(0xE3,0xF2,0xFD), bd=BIRU_MUDA, fc=BIRU_TUA)

# ── III-4: Setujui/Tolak RA ──────────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "III-03  |  Setujui atau Tolak Rencana Aksi",
    "Langkah persetujuan dan prosedur penolakan RA")
ftr(slide, 23, "—", "SIPITUNG — Bagian III: Kabag Umum | LLDIKTI Wilayah III")

rect(slide, 0.4, 1.35, 5.9, 5.75, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
rect(slide, 0.4, 1.35, 5.9, 0.45, fill=HIJAU_MID)
tb(slide, "Menyetujui Rencana Aksi", 0.55, 1.37, 5.65, 0.4, sz=Pt(12), bold=True, color=PUTIH)
for i, (num, ttl, dsc) in enumerate([
    ("1","Temukan Tim yang Perlu Disetujui","Cari baris dengan badge biru 'Menunggu Review'. Tombol Setujui ada di kolom Aksi."),
    ("2","Klik Tombol 'Setujui'","Klik tombol hijau 'Setujui'. Dialog konfirmasi akan muncul."),
    ("3","Konfirmasi","Klik 'Konfirmasi'. Status berubah menjadi 'Menunggu PPK' (badge amber)."),
    ("4","Verifikasi","Refresh halaman. Status Tim Kerja tersebut sudah berubah."),
    ("5","Ulangi untuk Tim Lain","Lakukan hal yang sama untuk setiap RA yang berstatus 'Menunggu Review'."),
]):
    step_row(slide, num, ttl, dsc, 0.4, 1.9+i*0.98, w=5.9, h=0.91, accent=HIJAU_MID)

note(slide, "ℹ  Setelah disetujui, dokumen diteruskan ke PPK. Kabag tidak perlu tindakan lanjutan.",
     0.4, 6.85, 5.9, 0.45, bg=RGBColor(0xE3,0xF2,0xFD), bd=BIRU_MUDA, fc=BIRU_TUA)

rect(slide, 6.7, 1.35, 6.25, 5.75, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
rect(slide, 6.7, 1.35, 6.25, 0.45, fill=MERAH)
tb(slide, "Menolak Rencana Aksi + Catatan", 6.85, 1.37, 6.0, 0.4, sz=Pt(12), bold=True, color=PUTIH)
for i, (num, ttl, dsc) in enumerate([
    ("1","Identifikasi Alasan","Tentukan poin spesifik yang perlu diperbaiki sebelum menolak."),
    ("2","Klik Tombol 'Tolak'","Klik tombol merah 'Tolak' pada kolom Aksi baris Tim Kerja."),
    ("3","Dialog Catatan Muncul","Kolom 'Catatan Revisi' wajib diisi sebelum dapat mengkonfirmasi."),
    ("4","Tulis Catatan Spesifik","Contoh: 'Target TW II IKU 2.1 terlalu rendah. Mohon disesuaikan dengan target tahunan.'"),
    ("5","Konfirmasi Tolak","Klik 'Konfirmasi Tolak'. Status berubah menjadi 'Ditolak' (badge merah)."),
]):
    step_row(slide, num, ttl, dsc, 6.7, 1.9+i*0.98, w=6.25, h=0.91, accent=MERAH)

note(slide, "⚠  Catatan revisi harus SPESIFIK dan JELAS. Catatan umum akan menyulitkan Tim Kerja dalam melakukan perbaikan.",
     6.7, 6.85, 6.25, 0.45)

# ── III-5: Review & Setujui/Tolak Laporan ────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "III-04  |  Review & Tindak Lanjut Laporan Pengukuran",
    "Cara mereview, menyetujui, atau mengembalikan laporan triwulan")
ftr(slide, 24, "—", "SIPITUNG — Bagian III: Kabag Umum | LLDIKTI Wilayah III")

ph(slide, "Tampilan Laporan Pengukuran Kabag", 0.4, 1.35, 4.8, 5.75)

tb(slide, "Review & Tindak Lanjut Laporan", 5.55, 1.38, 7.4, 0.38,
   sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.55, 1.38, 0.06, 0.38, fill=AKSEN)
for i, (num, ttl, dsc) in enumerate([
    ("1","Buka Pengukuran Kinerja","Klik 'Pengukuran Kinerja' di sidebar. Kartu laporan per Tim Kerja tampil di atas."),
    ("2","Pilih Periode","Pilih periode triwulan dari dropdown."),
    ("3","Identifikasi Menunggu Review","Cari kartu dengan badge biru 'Menunggu'. Berisi tombol Setujui dan Kembalikan."),
    ("4","Review Matriks Realisasi","Cek nilai realisasi, progress kegiatan, kendala, strategi tindak lanjut per IKU."),
    ("5a","Setujui — Klik 'Setujui'","Jika data lengkap: klik tombol hijau 'Setujui' → isi catatan (opsional) → Konfirmasi."),
    ("5b","Tolak — Klik 'Kembalikan'","Jika perlu perbaikan: klik merah 'Kembalikan' → isi rekomendasi (WAJIB) → Konfirmasi."),
]):
    acc = HIJAU_MID if "5a" in num else (MERAH if "5b" in num else HIJAU_KAB2)
    step_row(slide, num, ttl, dsc, 5.55, 1.9+i*0.83, w=7.4, h=0.76, accent=acc)

note(slide, "ℹ  Matriks realisasi menampilkan: Sasaran | IKU | PIC | Target PK | Target TW | Realisasi | Progress | Kendala | Strategi | Diisi Oleh.",
     5.55, 6.92, 7.4, 0.48, bg=RGBColor(0xE3,0xF2,0xFD), bd=BIRU_MUDA, fc=BIRU_TUA)

# ── III-6: Export PDF & Checklist ────────────────────────────────────────────
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
hdr(slide, "III-05  |  Export PDF & Checklist Review",
    "Mengekspor laporan ke PDF dan panduan checklist harian")
ftr(slide, 25, "—", "SIPITUNG — Bagian III: Kabag Umum | LLDIKTI Wilayah III")

rect(slide, 0.4, 1.35, 5.9, 3.8, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
rect(slide, 0.4, 1.35, 5.9, 0.42, fill=BIRU_MID)
tb(slide, "Export Laporan PDF", 0.55, 1.37, 5.65, 0.38, sz=Pt(12), bold=True, color=PUTIH)
for i, (num, ttl, dsc) in enumerate([
    ("1","Pilih Periode","Pada halaman Pengukuran Kinerja, pilih periode dari dropdown."),
    ("2","Klik Export PDF","Klik tombol 'Export PDF [TW]'. Halaman preview A3 landscape terbuka."),
    ("3","Cetak / Simpan PDF","Tekan Ctrl+P/Cmd+P → pilih 'Save as PDF' atau printer fisik."),
]):
    step_row(slide, num, ttl, dsc, 0.4, 1.87+i*0.98, w=5.9, h=0.9, accent=BIRU_MID)

note(slide, "ℹ  PDF berisi: Judul & Periode | Matriks Realisasi | Status per Tim Kerja | Area Tanda Tangan.",
     0.4, 5.25, 5.9, 0.45, bg=RGBColor(0xE3,0xF2,0xFD), bd=BIRU_MUDA, fc=BIRU_TUA)

# Checklist RA
rect(slide, 0.4, 5.85, 5.9, 1.62, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
rect(slide, 0.4, 5.85, 5.9, 0.38, fill=HIJAU_KAB2)
tb(slide, "Checklist Review RA", 0.55, 5.87, 5.65, 0.34, sz=Pt(10.5), bold=True, color=PUTIH)
for i, c in enumerate([
    "Semua IKU: target TW I–IV terisi & tidak ada nilai kosong",
    "Target triwulan kumulatif (TW IV = target tahunan)",
    "IKU kolaborasi ter-assign ke PIC primer dan co-PIC yang benar",
]):
    ty = 6.3+i*0.38
    tb(slide, "☐", 0.55, ty, 0.32, 0.34, sz=Pt(12), bold=True, color=HIJAU_MID, align=PP_ALIGN.CENTER)
    tb(slide, c, 0.92, ty+0.04, 5.28, 0.3, sz=Pt(9.5), color=ABU_TUA)

# Checklist Laporan
rect(slide, 6.7, 1.35, 6.25, 3.8, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
rect(slide, 6.7, 1.35, 6.25, 0.42, fill=HIJAU_TUA)
tb(slide, "Checklist Review Laporan Pengukuran", 6.85, 1.37, 6.0, 0.38, sz=Pt(12), bold=True, color=PUTIH)
for i, c in enumerate([
    "Semua IKU Tim Kerja sudah diisi nilai realisasi",
    "Nilai realisasi wajar dan konsisten dengan periode sebelumnya",
    "Kolom Progress Kegiatan terisi dengan deskripsi yang jelas",
    "Kolom Kendala diisi jika ada hambatan signifikan",
    "Kolom Strategi Tindak Lanjut diisi sesuai kendala",
    "IKU kolaborasi: kolom 'Diisi Oleh' menampilkan tim yang tepat",
]):
    ty = 1.87+i*0.55
    rect(slide, 6.7, ty, 6.25, 0.53, fill=ABU_MUDA if i%2==0 else PUTIH)
    tb(slide, "☐", 6.83, ty+0.06, 0.32, 0.4, sz=Pt(12), bold=True, color=HIJAU_MID, align=PP_ALIGN.CENTER)
    tb(slide, c, 7.2, ty+0.08, 5.65, 0.4, sz=Pt(9.5), color=ABU_TUA)

# FAQ Kabag
rect(slide, 6.7, 5.25, 6.25, 2.22, fill=PUTIH, line=RGBColor(0xCC,0xD9,0xEA), lw=Pt(1))
rect(slide, 6.7, 5.25, 6.25, 0.38, fill=MERAH)
tb(slide, "FAQ Singkat — Kabag Umum", 6.85, 5.27, 6.0, 0.34, sz=Pt(11), bold=True, color=PUTIH)
for i, (q, a) in enumerate([
    ("Tombol Setujui tidak muncul?", "Status dokumen bukan 'submitted'. Cek badge status Tim Kerja tersebut."),
    ("Bisakah saya batalkan persetujuan?", "Tidak. Hanya SuperAdmin yang bisa membuka kembali dokumen ppk_approved."),
    ("Laporan mana yang menunggu review?", "Cek Dashboard: angka 'Menunggu Review'. Klik untuk menuju halaman tersebut."),
]):
    ty = 5.72+i*0.58
    rect(slide, 6.7, ty, 6.25, 0.56, fill=ABU_MUDA if i%2==0 else PUTIH)
    tb(slide, f"Q: {q}", 6.85, ty+0.04, 6.0, 0.22, sz=Pt(9.5), bold=True, color=BIRU_TUA)
    tb(slide, f"A: {a}", 6.85, ty+0.28, 6.0, 0.24, sz=Pt(9), color=ABU_TUA)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE PENUTUP
# ════════════════════════════════════════════════════════════════════════════
slide = bs()
rect(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0x0A,0x14,0x2C))
rect(slide, 0, 3.5, 13.33, 0.08, fill=AKSEN)
tb(slide, "Terima Kasih", 1.0, 1.3, 11.33, 1.2,
   sz=Pt(44), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
tb(slide, "Buku Petunjuk Teknis SIPITUNG — LLDIKTI Wilayah III",
   1.0, 2.6, 11.33, 0.6, sz=Pt(16), color=AKSEN, align=PP_ALIGN.CENTER)
tb(slide, "Untuk pertanyaan dan bantuan teknis, hubungi Admin Sistem SIPITUNG",
   1.0, 3.75, 11.33, 0.5, sz=Pt(13), color=PUTIH, align=PP_ALIGN.CENTER)
tb(slide, "Versi 1.0  |  April 2026",
   1.0, 4.35, 11.33, 0.45, sz=Pt(11),
   color=RGBColor(0x88,0xA8,0xD0), align=PP_ALIGN.CENTER)
for i, (tag, col) in enumerate([
    ("Bagian I: Rencana Aksi",BIRU_MID),
    ("Bagian II: Pengukuran",HIJAU_MID),
    ("Bagian III: Kabag Umum",HIJAU_KAB2),
]):
    tx = 2.5+i*2.8
    rect(slide, tx, 5.15, 2.55, 0.58, fill=col)
    tb(slide, tag, tx+0.1, 5.22, 2.35, 0.44,
       sz=Pt(11), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SIMPAN
# ════════════════════════════════════════════════════════════════════════════
out = "Juknis_SIPITUNG_Lengkap.pptx"
prs.save(out)
total = len(prs.slides)
print(f"✓ File berhasil dibuat: {out}")
print(f"  Total slides: {total}")
print(f"\n  Struktur:")
print(f"    Slide  1      — Cover Utama")
print(f"    Slide  2      — Daftar Isi Master")
print(f"    Slide  3      — Divider Bagian I: Rencana Aksi")
print(f"    Slide  4–11   — Bagian I: Ketua Tim Kerja — Rencana Aksi")
print(f"    Slide 12      — Divider Bagian II: Pengukuran Kinerja")
print(f"    Slide 13–18   — Bagian II: Ketua Tim Kerja — Pengukuran Kinerja")
print(f"    Slide 19      — Divider Bagian III: Kabag Umum")
print(f"    Slide 20–25   — Bagian III: Kabag Umum")
print(f"    Slide {total}      — Penutup")
