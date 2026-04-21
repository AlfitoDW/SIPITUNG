"""
Script untuk generate Juknis SIPITUNG dalam format PowerPoint.
Jalankan: python3 generate_juknis.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import copy

# Warna tema LLDIKTI
BIRU_TUA = RGBColor(0x1A, 0x37, 0x6C)   # #1A376C
BIRU_MID = RGBColor(0x1E, 0x4D, 0x9B)   # #1E4D9B
BIRU_MUDA = RGBColor(0x2E, 0x75, 0xB6)  # #2E75B6
AKSEN = RGBColor(0xF5, 0xA6, 0x23)      # #F5A623 (kuning/emas)
PUTIH = RGBColor(0xFF, 0xFF, 0xFF)
ABU_MUDA = RGBColor(0xF2, 0xF4, 0xF8)   # background slide
ABU_TUA = RGBColor(0x44, 0x44, 0x44)
MERAH_PENTING = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank layout

# ──────────────────────────────────────────────────────────────────────────────
# HELPER FUNCTIONS
# ──────────────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h,
                 font_size=Pt(14), bold=False, color=ABU_TUA,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_header_bar(slide, title_text, subtitle_text=None):
    """Tambah header bar biru di atas slide."""
    add_rect(slide, 0, 0, 13.33, 1.1, fill_color=BIRU_TUA)
    add_rect(slide, 0, 1.1, 13.33, 0.06, fill_color=AKSEN)
    add_text_box(slide, title_text, 0.35, 0.12, 12.5, 0.65,
                 font_size=Pt(26), bold=True, color=PUTIH)
    if subtitle_text:
        add_text_box(slide, subtitle_text, 0.35, 0.7, 12.5, 0.38,
                     font_size=Pt(13), color=RGBColor(0xC8, 0xD8, 0xF0))


def add_footer(slide, page_num, total_pages):
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=BIRU_TUA)
    add_text_box(slide, "SIPITUNG — Petunjuk Teknis Ketua Tim Kerja | LLDIKTI Wilayah III",
                 0.35, 7.1, 10, 0.4, font_size=Pt(9), color=RGBColor(0xC8, 0xD8, 0xF0))
    add_text_box(slide, f"{page_num} / {total_pages}",
                 12.5, 7.1, 0.8, 0.4, font_size=Pt(9),
                 color=PUTIH, align=PP_ALIGN.RIGHT)


def add_step_box(slide, number, title, desc, l, t, w, h, accent=BIRU_MID):
    """Kotak langkah dengan nomor besar."""
    add_rect(slide, l, t, w, h, fill_color=PUTIH,
             line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(1))
    # Nomor
    num_w = 0.65
    add_rect(slide, l, t, num_w, h, fill_color=accent)
    add_text_box(slide, str(number), l + 0.08, t + 0.05, num_w - 0.1, h - 0.1,
                 font_size=Pt(28), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    # Konten
    add_text_box(slide, title, l + num_w + 0.1, t + 0.05, w - num_w - 0.15, 0.32,
                 font_size=Pt(12), bold=True, color=BIRU_TUA)
    add_text_box(slide, desc, l + num_w + 0.1, t + 0.35, w - num_w - 0.15, h - 0.42,
                 font_size=Pt(10), color=ABU_TUA, wrap=True)


def add_info_card(slide, icon_text, title, body, l, t, w, h, accent=BIRU_MUDA):
    add_rect(slide, l, t, w, h, fill_color=PUTIH,
             line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(1))
    add_rect(slide, l, t, w, 0.06, fill_color=accent)
    add_text_box(slide, icon_text + "  " + title,
                 l + 0.12, t + 0.12, w - 0.2, 0.32,
                 font_size=Pt(11), bold=True, color=BIRU_TUA)
    add_text_box(slide, body, l + 0.12, t + 0.46, w - 0.2, h - 0.55,
                 font_size=Pt(10), color=ABU_TUA, wrap=True)


TOTAL_SLIDES = 11

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# Background
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=ABU_MUDA)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=BIRU_TUA)

# Panel kiri dekoratif
add_rect(slide, 0, 0, 4.5, 7.5, fill_color=RGBColor(0x0F, 0x20, 0x44))
add_rect(slide, 4.5, 0, 0.08, 7.5, fill_color=AKSEN)

# Placeholder logo (kotak kosong dengan label)
add_rect(slide, 0.4, 0.4, 3.7, 1.6, fill_color=RGBColor(0x1A, 0x37, 0x6C),
         line_color=AKSEN, line_width=Pt(1.5))
add_text_box(slide, "[LOGO LLDIKTI\nWILAYAH III]",
             0.4, 0.4, 3.7, 1.6,
             font_size=Pt(11), color=RGBColor(0xC8, 0xD8, 0xF0),
             align=PP_ALIGN.CENTER, bold=False)

# Teks kiri
add_text_box(slide, "LLDIKTI WILAYAH III",
             0.3, 2.3, 3.9, 0.55,
             font_size=Pt(13), bold=True, color=AKSEN)
add_text_box(slide, "Kementerian Pendidikan Tinggi,\nSains dan Teknologi",
             0.3, 2.8, 3.9, 0.8,
             font_size=Pt(10), color=RGBColor(0xC8, 0xD8, 0xF0))

add_rect(slide, 0.3, 3.8, 3.9, 0.05, fill_color=AKSEN)

add_text_box(slide, "PETUNJUK TEKNIS",
             0.3, 4.0, 3.9, 0.45,
             font_size=Pt(11), color=RGBColor(0xC8, 0xD8, 0xF0), bold=True)
add_text_box(slide, "Penggunaan Aplikasi\nSIPITUNG",
             0.3, 4.4, 3.9, 0.9,
             font_size=Pt(14), color=PUTIH, bold=True)

add_text_box(slide, "Versi 1.0  |  April 2026",
             0.3, 6.9, 3.9, 0.4,
             font_size=Pt(9), color=RGBColor(0x88, 0xA8, 0xD0))

# Teks kanan — judul utama
add_text_box(slide, "SIPITUNG",
             4.9, 1.2, 8.0, 1.4,
             font_size=Pt(56), bold=True, color=PUTIH, align=PP_ALIGN.LEFT)

add_text_box(slide, "Sistem Informasi Perencanaan & Kinerja",
             4.9, 2.55, 8.0, 0.55,
             font_size=Pt(18), color=AKSEN, bold=False)

add_rect(slide, 4.9, 3.2, 8.0, 0.05, fill_color=RGBColor(0x88, 0xA8, 0xD0))

add_text_box(slide, "Modul Penyusunan Rencana Aksi",
             4.9, 3.35, 8.0, 0.5,
             font_size=Pt(16), color=PUTIH, bold=True)

add_text_box(slide,
             "Panduan Penggunaan untuk\nKetua Tim Kerja",
             4.9, 3.9, 8.0, 0.8,
             font_size=Pt(14), color=RGBColor(0xC8, 0xD8, 0xF0))

# Chips info
chip_data = [("Pengguna", "Ketua Tim Kerja"), ("Tahun", "2026"), ("Status", "Aktif")]
for i, (label, val) in enumerate(chip_data):
    cx = 4.9 + i * 2.7
    add_rect(slide, cx, 5.0, 2.4, 0.65, fill_color=RGBColor(0x1E, 0x4D, 0x9B),
             line_color=RGBColor(0x88, 0xA8, 0xD0), line_width=Pt(0.75))
    add_text_box(slide, label, cx + 0.1, 5.02, 2.2, 0.25,
                 font_size=Pt(8), color=AKSEN, bold=True)
    add_text_box(slide, val, cx + 0.1, 5.27, 2.2, 0.3,
                 font_size=Pt(11), color=PUTIH, bold=True)

# Nomor slide
add_text_box(slide, "1", 12.9, 7.2, 0.4, 0.3,
             font_size=Pt(8), color=RGBColor(0x88, 0xA8, 0xD0), align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — DAFTAR ISI
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=ABU_MUDA)
add_header_bar(slide, "Daftar Isi", "Petunjuk Teknis SIPITUNG — Ketua Tim Kerja")
add_footer(slide, 2, TOTAL_SLIDES)

toc_items = [
    ("01", "Pendahuluan & Tujuan", "Latar belakang dan maksud penggunaan juknis ini"),
    ("02", "Ruang Lingkup", "Batasan pengguna dan fitur yang dibahas"),
    ("03", "Ketentuan Umum", "Aturan dan kebijakan yang perlu diketahui"),
    ("04", "Login ke Aplikasi", "Cara masuk ke SIPITUNG"),
    ("05", "Navigasi Menu Perencanaan", "Akses menu Rencana Aksi > Penyusunan"),
    ("06", "Mengisi Target Triwulan", "Cara edit IKU dan mengisi target TW I–IV"),
    ("07", "Menambah Rencana Kegiatan", "Cara mengisi kegiatan per triwulan"),
    ("08", "Submit ke Kabag Umum", "Cara mengirim dokumen untuk persetujuan"),
    ("09", "Revisi & Submit Ulang", "Prosedur jika dokumen dikembalikan"),
    ("10", "FAQ", "Pertanyaan yang sering diajukan"),
]

cols = 2
rows = 5
for i, (num, title, desc) in enumerate(toc_items):
    col = i % cols
    row = i // cols
    lx = 0.4 + col * 6.5
    ty = 1.35 + row * 1.12
    add_rect(slide, lx, ty, 6.1, 1.0, fill_color=PUTIH,
             line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(1))
    add_rect(slide, lx, ty, 0.7, 1.0, fill_color=BIRU_MID)
    add_text_box(slide, num, lx + 0.05, ty + 0.2, 0.6, 0.55,
                 font_size=Pt(16), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, lx + 0.78, ty + 0.06, 5.2, 0.32,
                 font_size=Pt(11), bold=True, color=BIRU_TUA)
    add_text_box(slide, desc, lx + 0.78, ty + 0.38, 5.2, 0.55,
                 font_size=Pt(9), color=ABU_TUA, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PENDAHULUAN & TUJUAN
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=ABU_MUDA)
add_header_bar(slide, "01  |  Pendahuluan & Tujuan", "Latar belakang dan maksud dokumen ini")
add_footer(slide, 3, TOTAL_SLIDES)

add_text_box(slide, "Latar Belakang",
             0.4, 1.35, 6.1, 0.38, font_size=Pt(13), bold=True, color=BIRU_TUA)
add_rect(slide, 0.4, 1.35, 0.06, 0.38, fill_color=AKSEN)
add_text_box(slide,
    "SIPITUNG (Sistem Informasi Perencanaan & Kinerja) dikembangkan untuk mendukung "
    "proses perencanaan kerja di lingkungan LLDIKTI Wilayah III secara digital, "
    "transparan, dan terstruktur. Setiap Tim Kerja diwajibkan menyusun Rencana Aksi "
    "yang memuat target triwulan dan rencana kegiatan untuk setiap Indikator Kinerja "
    "Utama (IKU) yang menjadi tanggung jawab mereka.",
    0.5, 1.8, 6.0, 1.8, font_size=Pt(10.5), color=ABU_TUA, wrap=True)

add_text_box(slide, "Tujuan Juknis",
             0.4, 3.75, 6.1, 0.38, font_size=Pt(13), bold=True, color=BIRU_TUA)
add_rect(slide, 0.4, 3.75, 0.06, 0.38, fill_color=AKSEN)

tujuan = [
    "Memberikan panduan penggunaan modul Penyusunan Rencana Aksi bagi Ketua Tim Kerja.",
    "Memastikan data yang diinput sesuai format dan ketentuan yang berlaku.",
    "Meminimalkan kesalahan dalam proses input dan pengajuan Rencana Aksi.",
    "Menjadi referensi saat terjadi kendala teknis di lapangan.",
]
for i, t in enumerate(tujuan):
    ty = 4.2 + i * 0.52
    add_rect(slide, 0.5, ty + 0.09, 0.22, 0.22, fill_color=BIRU_MID)
    add_text_box(slide, str(i + 1), 0.5, ty + 0.05, 0.22, 0.28,
                 font_size=Pt(9), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    add_text_box(slide, t, 0.8, ty, 5.8, 0.5,
                 font_size=Pt(10.5), color=ABU_TUA, wrap=True)

# Panel kanan
add_rect(slide, 7.0, 1.25, 6.0, 5.8, fill_color=PUTIH,
         line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(1))
add_rect(slide, 7.0, 1.25, 6.0, 0.08, fill_color=BIRU_MUDA)

add_text_box(slide, "Informasi Dokumen",
             7.2, 1.38, 5.6, 0.38, font_size=Pt(12), bold=True, color=BIRU_TUA)

doc_info = [
    ("Nama Sistem", "SIPITUNG"),
    ("Versi", "1.0"),
    ("Tanggal Terbit", "April 2026"),
    ("Pengguna", "Ketua Tim Kerja"),
    ("Modul", "Perencanaan > Rencana Aksi > Penyusunan"),
    ("Instansi", "LLDIKTI Wilayah III"),
    ("Kontak Dukungan", "Admin Sistem — ext. 100"),
]
for i, (k, v) in enumerate(doc_info):
    ty = 1.9 + i * 0.68
    add_rect(slide, 7.0, ty, 6.0, 0.68,
             fill_color=ABU_MUDA if i % 2 == 0 else PUTIH)
    add_text_box(slide, k, 7.15, ty + 0.08, 2.4, 0.52,
                 font_size=Pt(10), bold=True, color=BIRU_TUA)
    add_text_box(slide, v, 9.6, ty + 0.08, 3.3, 0.52,
                 font_size=Pt(10), color=ABU_TUA)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — RUANG LINGKUP & KETENTUAN UMUM
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=ABU_MUDA)
add_header_bar(slide, "02–03  |  Ruang Lingkup & Ketentuan Umum")
add_footer(slide, 4, TOTAL_SLIDES)

# Ruang Lingkup (kiri)
add_rect(slide, 0.4, 1.35, 5.8, 5.75, fill_color=PUTIH,
         line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(1))
add_rect(slide, 0.4, 1.35, 5.8, 0.45, fill_color=BIRU_MID)
add_text_box(slide, "02  |  Ruang Lingkup",
             0.55, 1.38, 5.5, 0.4, font_size=Pt(12), bold=True, color=PUTIH)

rl_items = [
    ("Pengguna", "Ketua Tim Kerja yang telah memiliki akun aktif di SIPITUNG (username: ketua.xxx)"),
    ("Modul", "Perencanaan > Rencana Aksi > Penyusunan"),
    ("Periode", "Setiap tahun anggaran, mengikuti siklus perencanaan LLDIKTI Wilayah III"),
    ("Cakupan Data", "IKU mandiri (milik tim sendiri) dan IKU kolaborasi (bersama tim lain)"),
    ("Di Luar Lingkup", "Pengaturan akun, manajemen IKU, dan persetujuan dokumen (kewenangan Kabag Umum)"),
]
for i, (k, v) in enumerate(rl_items):
    ty = 1.9 + i * 0.98
    add_rect(slide, 0.4, ty, 5.8, 0.95,
             fill_color=ABU_MUDA if i % 2 == 0 else PUTIH)
    add_text_box(slide, k, 0.55, ty + 0.06, 1.5, 0.32,
                 font_size=Pt(9), bold=True, color=BIRU_MID)
    add_text_box(slide, v, 0.55, ty + 0.36, 5.5, 0.55,
                 font_size=Pt(9.5), color=ABU_TUA, wrap=True)

# Ketentuan Umum (kanan)
add_rect(slide, 6.9, 1.35, 6.05, 5.75, fill_color=PUTIH,
         line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(1))
add_rect(slide, 6.9, 1.35, 6.05, 0.45, fill_color=BIRU_TUA)
add_text_box(slide, "03  |  Ketentuan Umum",
             7.05, 1.38, 5.8, 0.4, font_size=Pt(12), bold=True, color=PUTIH)

ku_items = [
    ("1", "Akun Aktif", "Pengguna harus memiliki akun dengan status aktif. Hubungi Admin jika tidak dapat login."),
    ("2", "Kelengkapan Data", "Semua IKU yang tampil wajib diisi target TW I–IV sebelum dokumen dapat disubmit."),
    ("3", "Dokumen Tersubmit", "Dokumen yang telah disubmit tidak dapat diubah, kecuali dikembalikan oleh Kabag Umum dengan catatan revisi."),
    ("4", "IKU Kolaborasi", "Target pada IKU kolaborasi diisi bersama, namun masing-masing Tim Kerja mengisi rencana kegiatannya sendiri."),
    ("5", "Kerahasiaan Akun", "Username dan password bersifat rahasia. Jangan berbagi akun dengan pihak lain."),
    ("6", "Bantuan Teknis", "Jika mengalami kendala teknis, hubungi Admin Sistem melalui kontak yang tersedia."),
]
for i, (num, title, desc) in enumerate(ku_items):
    ty = 1.88 + i * 0.84
    add_rect(slide, 6.9, ty, 6.05, 0.82,
             fill_color=ABU_MUDA if i % 2 == 0 else PUTIH)
    add_rect(slide, 6.9, ty, 0.4, 0.82, fill_color=BIRU_MUDA if i % 2 == 0 else BIRU_MID)
    add_text_box(slide, num, 6.9, ty + 0.18, 0.4, 0.4,
                 font_size=Pt(13), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, 7.38, ty + 0.04, 5.4, 0.28,
                 font_size=Pt(10), bold=True, color=BIRU_TUA)
    add_text_box(slide, desc, 7.38, ty + 0.32, 5.4, 0.46,
                 font_size=Pt(9), color=ABU_TUA, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — LANGKAH 1: LOGIN
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=ABU_MUDA)
add_header_bar(slide, "04  |  Login ke Aplikasi SIPITUNG", "Langkah pertama sebelum menggunakan sistem")
add_footer(slide, 5, TOTAL_SLIDES)

# Panel screenshot placeholder (kiri)
add_rect(slide, 0.4, 1.35, 5.5, 5.75, fill_color=RGBColor(0xE8, 0xEE, 0xF8),
         line_color=BIRU_MUDA, line_width=Pt(1.5))
add_rect(slide, 0.4, 1.35, 5.5, 0.38, fill_color=BIRU_MUDA)
add_text_box(slide, "Tampilan Halaman Login",
             0.55, 1.37, 5.2, 0.35, font_size=Pt(10), bold=True, color=PUTIH)
add_text_box(slide,
    "[ SCREENSHOT HALAMAN LOGIN ]\n\nTambahkan screenshot tampilan\nhalaman login aplikasi SIPITUNG\ndi sini",
    0.9, 2.4, 4.6, 4.0,
    font_size=Pt(11), color=BIRU_MUDA, align=PP_ALIGN.CENTER, italic=True)
add_rect(slide, 1.5, 3.5, 3.3, 2.0, fill_color=RGBColor(0xD0, 0xDC, 0xF0),
         line_color=BIRU_MUDA, line_width=Pt(0.75))
add_text_box(slide, "Username\n_______________\n\nPassword\n_______________\n\n[ LOGIN ]",
             1.6, 3.55, 3.1, 1.9, font_size=Pt(9), color=BIRU_TUA, align=PP_ALIGN.CENTER)

# Langkah-langkah (kanan)
add_text_box(slide, "Langkah-Langkah Login",
             6.3, 1.35, 6.6, 0.4, font_size=Pt(13), bold=True, color=BIRU_TUA)
add_rect(slide, 6.3, 1.35, 0.06, 0.4, fill_color=AKSEN)

steps = [
    ("1", "Buka Aplikasi", 'Buka browser (Chrome/Firefox) dan akses URL SIPITUNG yang diberikan Admin Sistem.'),
    ("2", "Masukkan Username", 'Ketik username Anda (format: ketua.xxx, contoh: ketua.keuangan) pada kolom Username.'),
    ("3", "Masukkan Password", 'Ketik password Anda pada kolom Password. Password bersifat case-sensitive.'),
    ("4", "Klik Tombol Login", 'Klik tombol "Masuk" atau tekan Enter pada keyboard.'),
    ("5", "Verifikasi Berhasil", 'Jika berhasil, Anda akan diarahkan ke Dashboard Ketua Tim Kerja.'),
]
for i, (num, title, desc) in enumerate(steps):
    ty = 1.9 + i * 1.02
    add_rect(slide, 6.3, ty, 6.6, 0.95, fill_color=PUTIH,
             line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(0.75))
    add_rect(slide, 6.3, ty, 0.55, 0.95, fill_color=BIRU_MID)
    add_text_box(slide, num, 6.3, ty + 0.2, 0.55, 0.55,
                 font_size=Pt(20), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, 6.93, ty + 0.05, 5.9, 0.3,
                 font_size=Pt(10.5), bold=True, color=BIRU_TUA)
    add_text_box(slide, desc, 6.93, ty + 0.35, 5.9, 0.55,
                 font_size=Pt(9.5), color=ABU_TUA, wrap=True)

# Catatan penting
add_rect(slide, 6.3, 7.0, 6.6, 0.0, fill_color=PUTIH)  # spacer
note_box = add_rect(slide, 6.3, 6.95, 6.6, 0.6, fill_color=RGBColor(0xFF, 0xF3, 0xCD),
                    line_color=RGBColor(0xF5, 0xA6, 0x23), line_width=Pt(1))
add_text_box(slide, "⚠  Perhatian: Jika login gagal, periksa kembali username dan password. Pastikan Caps Lock tidak aktif. Hubungi Admin jika akun terkunci.",
             6.45, 6.97, 6.3, 0.56, font_size=Pt(9), color=RGBColor(0x7D, 0x60, 0x08), wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LANGKAH 2: NAVIGASI MENU
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=ABU_MUDA)
add_header_bar(slide, "05  |  Navigasi Menu Penyusunan Rencana Aksi", "Menemukan dan membuka halaman yang tepat")
add_footer(slide, 6, TOTAL_SLIDES)

# Breadcrumb visual
add_rect(slide, 0.4, 1.35, 12.5, 0.55, fill_color=PUTIH,
         line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(1))
breadcrumbs = ["Dashboard", "Perencanaan", "Rencana Aksi", "Penyusunan"]
bw = 2.8
for i, bc in enumerate(breadcrumbs):
    bx = 0.5 + i * bw
    bg = BIRU_TUA if i == len(breadcrumbs) - 1 else (BIRU_MUDA if i > 0 else RGBColor(0xE0, 0xE8, 0xF5))
    fc = PUTIH if i >= 1 else BIRU_TUA
    add_rect(slide, bx, 1.38, bw - 0.1, 0.49, fill_color=bg)
    add_text_box(slide, ("▶  " if i > 0 else "") + bc,
                 bx + 0.1, 1.41, bw - 0.15, 0.43,
                 font_size=Pt(10), bold=(i == len(breadcrumbs)-1), color=fc)

# Screenshot placeholder
add_rect(slide, 0.4, 2.05, 4.5, 4.95, fill_color=RGBColor(0xE8, 0xEE, 0xF8),
         line_color=BIRU_MUDA, line_width=Pt(1.5))
add_rect(slide, 0.4, 2.05, 4.5, 0.38, fill_color=BIRU_MUDA)
add_text_box(slide, "Tampilan Menu Sidebar",
             0.55, 2.07, 4.2, 0.35, font_size=Pt(10), bold=True, color=PUTIH)
add_text_box(slide,
    "[ SCREENSHOT SIDEBAR MENU ]\n\nTambahkan screenshot menu\nnavigasi sidebar di sini\n\nTunjukkan:\n• Menu 'Perencanaan'\n• Submenu 'Rencana Aksi'\n• Submenu 'Penyusunan'",
    0.6, 2.7, 4.1, 4.1,
    font_size=Pt(10), color=BIRU_MUDA, align=PP_ALIGN.CENTER, italic=True)

# Langkah navigasi
add_text_box(slide, "Langkah Navigasi",
             5.25, 2.05, 7.7, 0.38, font_size=Pt(13), bold=True, color=BIRU_TUA)
add_rect(slide, 5.25, 2.05, 0.06, 0.38, fill_color=AKSEN)

nav_steps = [
    ("1", "Buka Sidebar", "Setelah login, sidebar navigasi akan tampil di sisi kiri layar."),
    ("2", "Klik 'Perencanaan'", "Klik menu 'Perencanaan' pada sidebar untuk memperluas submenu."),
    ("3", "Pilih 'Rencana Aksi'", "Dari submenu yang muncul, klik 'Rencana Aksi'."),
    ("4", "Pilih 'Penyusunan'", "Klik submenu 'Penyusunan' untuk masuk ke halaman daftar IKU tim Anda."),
    ("5", "Verifikasi Halaman", "Pastikan judul halaman menampilkan nama Tim Kerja Anda dan daftar IKU tampil dengan benar."),
]
for i, (num, title, desc) in enumerate(nav_steps):
    ty = 2.55 + i * 0.9
    add_rect(slide, 5.25, ty, 7.7, 0.83, fill_color=PUTIH,
             line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(0.75))
    add_rect(slide, 5.25, ty, 0.5, 0.83, fill_color=BIRU_MID)
    add_text_box(slide, num, 5.25, ty + 0.17, 0.5, 0.5,
                 font_size=Pt(18), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, 5.82, ty + 0.04, 7.05, 0.28,
                 font_size=Pt(10.5), bold=True, color=BIRU_TUA)
    add_text_box(slide, desc, 5.82, ty + 0.34, 7.05, 0.45,
                 font_size=Pt(9.5), color=ABU_TUA, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — LANGKAH 3: MENGISI TARGET TRIWULAN
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=ABU_MUDA)
add_header_bar(slide, "06  |  Mengisi Target Triwulan (TW I–IV)", "Klik ikon pensil pada setiap IKU untuk membuka form isian")
add_footer(slide, 7, TOTAL_SLIDES)

# Screenshot placeholder
add_rect(slide, 0.4, 1.35, 5.3, 5.75, fill_color=RGBColor(0xE8, 0xEE, 0xF8),
         line_color=BIRU_MUDA, line_width=Pt(1.5))
add_rect(slide, 0.4, 1.35, 5.3, 0.38, fill_color=BIRU_MUDA)
add_text_box(slide, "Tampilan Form Isian Target",
             0.55, 1.37, 5.0, 0.35, font_size=Pt(10), bold=True, color=PUTIH)
add_text_box(slide,
    "[ SCREENSHOT FORM EDIT TARGET ]\n\nTambahkan screenshot tampilan\nform isian target triwulan\n\nTunjukkan:\n• Daftar IKU\n• Tombol ikon pensil (edit)\n• Kolom TW I, TW II, TW III, TW IV\n• Tombol Simpan",
    0.65, 2.1, 4.8, 4.7,
    font_size=Pt(10), color=BIRU_MUDA, align=PP_ALIGN.CENTER, italic=True)

# Konten kanan
add_text_box(slide, "Cara Mengisi Target Triwulan",
             6.05, 1.35, 7.0, 0.38, font_size=Pt(13), bold=True, color=BIRU_TUA)
add_rect(slide, 6.05, 1.35, 0.06, 0.38, fill_color=AKSEN)

target_steps = [
    ("1", "Temukan IKU di Daftar", "Pada halaman Penyusunan, Anda akan melihat daftar IKU mandiri dan kolaborasi yang menjadi tanggung jawab tim Anda."),
    ("2", "Klik Ikon Pensil (✏)", "Klik ikon pensil (edit) pada baris IKU yang ingin diisi. Sebuah form isian akan terbuka."),
    ("3", "Isi Target TW I–IV", "Isi nilai target untuk masing-masing triwulan:\n• TW I : Target akhir Maret\n• TW II: Target akhir Juni\n• TW III: Target akhir September\n• TW IV: Target akhir Desember"),
    ("4", "Klik Simpan", 'Klik tombol "Simpan" untuk menyimpan data. Status IKU akan berubah menandai telah diisi.'),
    ("5", "Ulangi untuk Semua IKU", "Lakukan langkah 2–4 untuk setiap IKU yang tampil di daftar."),
]
for i, (num, title, desc) in enumerate(target_steps):
    ty = 1.9 + i * 0.97
    add_rect(slide, 6.05, ty, 7.0, 0.9, fill_color=PUTIH,
             line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(0.75))
    add_rect(slide, 6.05, ty, 0.5, 0.9, fill_color=BIRU_MID)
    add_text_box(slide, num, 6.05, ty + 0.17, 0.5, 0.55,
                 font_size=Pt(18), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, 6.63, ty + 0.04, 6.35, 0.28,
                 font_size=Pt(10.5), bold=True, color=BIRU_TUA)
    add_text_box(slide, desc, 6.63, ty + 0.34, 6.35, 0.52,
                 font_size=Pt(9.5), color=ABU_TUA, wrap=True)

# Catatan
add_rect(slide, 6.05, 6.75, 7.0, 0.6, fill_color=RGBColor(0xE8, 0xF5, 0xE9),
         line_color=RGBColor(0x27, 0xAE, 0x60), line_width=Pt(1))
add_text_box(slide, "✓  Tips: Target triwulan bersifat kumulatif. Pastikan nilai TW IV sesuai dengan target tahunan yang telah ditetapkan dalam dokumen PK Awal.",
             6.2, 6.77, 6.7, 0.56, font_size=Pt(9), color=RGBColor(0x1A, 0x6B, 0x30), wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — LANGKAH 4: MENAMBAH RENCANA KEGIATAN
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=ABU_MUDA)
add_header_bar(slide, "07  |  Menambahkan Rencana Kegiatan", "Klik ikon list pada setiap IKU untuk mengelola kegiatan")
add_footer(slide, 8, TOTAL_SLIDES)

# Screenshot placeholder
add_rect(slide, 0.4, 1.35, 5.3, 5.75, fill_color=RGBColor(0xE8, 0xEE, 0xF8),
         line_color=BIRU_MUDA, line_width=Pt(1.5))
add_rect(slide, 0.4, 1.35, 5.3, 0.38, fill_color=BIRU_MUDA)
add_text_box(slide, "Tampilan Form Rencana Kegiatan",
             0.55, 1.37, 5.0, 0.35, font_size=Pt(10), bold=True, color=PUTIH)
add_text_box(slide,
    "[ SCREENSHOT FORM KEGIATAN ]\n\nTambahkan screenshot tampilan\nform pengelolaan rencana\nkegiatan\n\nTunjukkan:\n• Tombol ikon list (kegiatan)\n• Form tambah kegiatan\n• Kolom nama kegiatan\n• Pilihan triwulan",
    0.65, 2.1, 4.8, 4.7,
    font_size=Pt(10), color=BIRU_MUDA, align=PP_ALIGN.CENTER, italic=True)

# Konten kanan
add_text_box(slide, "Cara Menambahkan Rencana Kegiatan",
             6.05, 1.35, 7.0, 0.38, font_size=Pt(13), bold=True, color=BIRU_TUA)
add_rect(slide, 6.05, 1.35, 0.06, 0.38, fill_color=AKSEN)

kegiatan_steps = [
    ("1", "Klik Ikon Kegiatan (☰)", "Pada baris IKU, klik ikon list/kegiatan untuk membuka panel rencana kegiatan IKU tersebut."),
    ("2", "Klik 'Tambah Kegiatan'", "Klik tombol Tambah Kegiatan untuk membuka form isian kegiatan baru."),
    ("3", "Isi Nama Kegiatan", "Tuliskan nama kegiatan secara jelas dan spesifik (contoh: 'Rapat Koordinasi Penyusunan Anggaran TW I')."),
    ("4", "Pilih Triwulan", "Pilih triwulan pelaksanaan kegiatan (TW I, TW II, TW III, atau TW IV)."),
    ("5", "Simpan Kegiatan", 'Klik "Simpan". Kegiatan akan muncul dalam daftar rencana kegiatan IKU tersebut.'),
    ("6", "Ulangi Sesuai Kebutuhan", "Tambahkan semua kegiatan yang direncanakan. Setiap IKU dapat memiliki beberapa kegiatan di tiap triwulan."),
]
for i, (num, title, desc) in enumerate(kegiatan_steps):
    ty = 1.9 + i * 0.82
    add_rect(slide, 6.05, ty, 7.0, 0.76, fill_color=PUTIH,
             line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(0.75))
    add_rect(slide, 6.05, ty, 0.45, 0.76, fill_color=BIRU_MID)
    add_text_box(slide, num, 6.05, ty + 0.14, 0.45, 0.48,
                 font_size=Pt(16), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, 6.57, ty + 0.03, 6.4, 0.28,
                 font_size=Pt(10.5), bold=True, color=BIRU_TUA)
    add_text_box(slide, desc, 6.57, ty + 0.32, 6.4, 0.42,
                 font_size=Pt(9.5), color=ABU_TUA, wrap=True)

# Catatan
add_rect(slide, 6.05, 6.82, 7.0, 0.6, fill_color=RGBColor(0xFF, 0xF3, 0xCD),
         line_color=AKSEN, line_width=Pt(1))
add_text_box(slide, "⚠  Perhatian: Minimal satu kegiatan harus ditambahkan per IKU sebelum dokumen dapat disubmit. Pastikan nama kegiatan deskriptif dan mudah dipahami.",
             6.2, 6.84, 6.7, 0.56, font_size=Pt(9), color=RGBColor(0x7D, 0x60, 0x08), wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — LANGKAH 5: SUBMIT KE KABAG
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=ABU_MUDA)
add_header_bar(slide, "08  |  Submit Rencana Aksi ke Kabag Umum", "Langkah pengajuan setelah semua IKU terisi lengkap")
add_footer(slide, 9, TOTAL_SLIDES)

# Alur persetujuan (atas)
add_rect(slide, 0.4, 1.35, 12.5, 1.5, fill_color=PUTIH,
         line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(1))
alur = [
    ("Pengisian\nData", BIRU_MUDA, "Ketua Tim\nKerja"),
    ("Submit ke\nKabag", BIRU_MID, "Ketua Tim\nKerja"),
    ("Review &\nPersetujuan", BIRU_TUA, "Kabag\nUmum"),
    ("Dokumen\nDisetujui", RGBColor(0x27, 0xAE, 0x60), "Kabag\nUmum"),
]
for i, (step, color, actor) in enumerate(alur):
    ax = 0.7 + i * 3.1
    add_rect(slide, ax, 1.45, 2.7, 1.3, fill_color=color)
    add_text_box(slide, step, ax + 0.1, 1.5, 2.5, 0.65,
                 font_size=Pt(10), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    add_text_box(slide, actor, ax + 0.1, 2.1, 2.5, 0.6,
                 font_size=Pt(8.5), color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
    if i < len(alur) - 1:
        add_text_box(slide, "▶", ax + 2.72, 1.8, 0.35, 0.55,
                     font_size=Pt(16), color=BIRU_MUDA, align=PP_ALIGN.CENTER)

# Screenshot placeholder
add_rect(slide, 0.4, 3.0, 5.3, 4.1, fill_color=RGBColor(0xE8, 0xEE, 0xF8),
         line_color=BIRU_MUDA, line_width=Pt(1.5))
add_rect(slide, 0.4, 3.0, 5.3, 0.38, fill_color=BIRU_MUDA)
add_text_box(slide, "Tampilan Tombol Submit",
             0.55, 3.02, 5.0, 0.35, font_size=Pt(10), bold=True, color=PUTIH)
add_text_box(slide,
    "[ SCREENSHOT TOMBOL SUBMIT ]\n\nTambahkan screenshot tombol\nSubmit ke Kabag Umum\n\nTunjukkan:\n• Status kelengkapan IKU\n• Tombol Submit\n• Dialog konfirmasi",
    0.65, 3.65, 4.8, 3.2,
    font_size=Pt(10), color=BIRU_MUDA, align=PP_ALIGN.CENTER, italic=True)

# Langkah submit
add_text_box(slide, "Cara Submit Dokumen",
             6.05, 3.0, 7.0, 0.38, font_size=Pt(13), bold=True, color=BIRU_TUA)
add_rect(slide, 6.05, 3.0, 0.06, 0.38, fill_color=AKSEN)

submit_steps = [
    ("1", "Periksa Kelengkapan", "Pastikan semua IKU telah diisi target TW I–IV dan minimal satu rencana kegiatan."),
    ("2", "Klik Tombol Submit", 'Klik tombol "Submit ke Kabag Umum" yang tersedia di halaman Penyusunan.'),
    ("3", "Konfirmasi Submit", "Sebuah dialog konfirmasi akan muncul. Baca dengan seksama sebelum mengklik Konfirmasi."),
    ("4", "Tunggu Notifikasi", "Setelah submit berhasil, status dokumen berubah menjadi 'Menunggu Persetujuan'."),
]
for i, (num, title, desc) in enumerate(submit_steps):
    ty = 3.5 + i * 0.85
    add_rect(slide, 6.05, ty, 7.0, 0.78, fill_color=PUTIH,
             line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(0.75))
    add_rect(slide, 6.05, ty, 0.5, 0.78, fill_color=BIRU_TUA)
    add_text_box(slide, num, 6.05, ty + 0.15, 0.5, 0.48,
                 font_size=Pt(18), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, 6.63, ty + 0.04, 6.35, 0.28,
                 font_size=Pt(10.5), bold=True, color=BIRU_TUA)
    add_text_box(slide, desc, 6.63, ty + 0.34, 6.35, 0.4,
                 font_size=Pt(9.5), color=ABU_TUA, wrap=True)

# Catatan kritis
add_rect(slide, 6.05, 6.9, 7.0, 0.55, fill_color=RGBColor(0xFD, 0xED, 0xEC),
         line_color=MERAH_PENTING, line_width=Pt(1))
add_text_box(slide, "⛔  Penting: Dokumen yang telah disubmit TIDAK DAPAT diedit. Pastikan semua data sudah benar sebelum menekan tombol Submit.",
             6.2, 6.92, 6.7, 0.51, font_size=Pt(9), color=MERAH_PENTING, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — LANGKAH 6: REVISI & SUBMIT ULANG
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=ABU_MUDA)
add_header_bar(slide, "09  |  Revisi & Submit Ulang", "Prosedur jika dokumen dikembalikan oleh Kabag Umum")
add_footer(slide, 10, TOTAL_SLIDES)

# Flowchart revisi
add_rect(slide, 0.4, 1.35, 12.5, 0.55, fill_color=PUTIH,
         line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(1))
add_text_box(slide, "Alur Revisi Dokumen Rencana Aksi",
             0.55, 1.38, 12.2, 0.45, font_size=Pt(11), bold=True, color=BIRU_TUA)

flow_items = [
    ("Dikembalikan\noleh Kabag", MERAH_PENTING),
    ("Baca\nCatatan\nRevisi", BIRU_MID),
    ("Perbaiki\nData IKU", BIRU_MUDA),
    ("Submit\nUlang", RGBColor(0x27, 0xAE, 0x60)),
]
for i, (label, color) in enumerate(flow_items):
    fx = 0.6 + i * 3.1
    add_rect(slide, fx, 2.05, 2.7, 1.1, fill_color=color)
    add_text_box(slide, label, fx + 0.1, 2.15, 2.5, 0.9,
                 font_size=Pt(11), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        add_text_box(slide, "▶", fx + 2.73, 2.45, 0.35, 0.45,
                     font_size=Pt(14), color=BIRU_MID, align=PP_ALIGN.CENTER)

# Kolom kiri: cara cek status
add_rect(slide, 0.4, 3.35, 6.0, 3.7, fill_color=PUTIH,
         line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(1))
add_rect(slide, 0.4, 3.35, 6.0, 0.42, fill_color=BIRU_MID)
add_text_box(slide, "Cara Memeriksa Status & Catatan Revisi",
             0.55, 3.37, 5.7, 0.38, font_size=Pt(11), bold=True, color=PUTIH)

status_steps = [
    "Buka halaman Perencanaan > Rencana Aksi > Penyusunan.",
    "Perhatikan status dokumen di bagian atas halaman.",
    "Jika status 'Dikembalikan', klik tombol Lihat Catatan untuk membaca keterangan dari Kabag Umum.",
    "Catat semua poin yang perlu diperbaiki sebelum melakukan perubahan.",
]
for i, s in enumerate(status_steps):
    ty = 3.88 + i * 0.73
    add_text_box(slide, f"0{i+1}", 0.55, ty + 0.05, 0.4, 0.5,
                 font_size=Pt(14), bold=True, color=BIRU_MID)
    add_text_box(slide, s, 1.0, ty, 5.3, 0.7,
                 font_size=Pt(10), color=ABU_TUA, wrap=True)

# Kolom kanan: langkah submit ulang
add_rect(slide, 6.9, 3.35, 6.05, 3.7, fill_color=PUTIH,
         line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(1))
add_rect(slide, 6.9, 3.35, 6.05, 0.42, fill_color=RGBColor(0x27, 0xAE, 0x60))
add_text_box(slide, "Langkah Perbaikan & Submit Ulang",
             7.05, 3.37, 5.7, 0.38, font_size=Pt(11), bold=True, color=PUTIH)

revisi_steps = [
    "Klik ikon pensil (✏) pada IKU yang perlu diperbaiki.",
    "Lakukan perubahan sesuai catatan revisi dari Kabag Umum.",
    "Klik Simpan setelah setiap perubahan.",
    'Setelah semua perbaikan selesai, klik kembali tombol "Submit ke Kabag Umum".',
]
for i, s in enumerate(revisi_steps):
    ty = 3.88 + i * 0.73
    add_rect(slide, 6.9, ty, 0.45, 0.65, fill_color=RGBColor(0x27, 0xAE, 0x60))
    add_text_box(slide, str(i + 1), 6.9, ty + 0.1, 0.45, 0.45,
                 font_size=Pt(14), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    add_text_box(slide, s, 7.43, ty, 5.4, 0.7,
                 font_size=Pt(10), color=ABU_TUA, wrap=True)

# Catatan
add_rect(slide, 0.4, 7.0, 12.5, 0.55, fill_color=RGBColor(0xFF, 0xF3, 0xCD),
         line_color=AKSEN, line_width=Pt(1))
add_text_box(slide, "⚠  Tidak ada batasan jumlah submit ulang. Namun, pastikan semua catatan revisi sudah ditangani agar persetujuan dapat diberikan.",
             0.6, 7.02, 12.2, 0.51, font_size=Pt(9.5), color=RGBColor(0x7D, 0x60, 0x08), wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — FAQ & PENUTUP
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=ABU_MUDA)
add_header_bar(slide, "10  |  FAQ & Penutup", "Pertanyaan yang Sering Diajukan")
add_footer(slide, 11, TOTAL_SLIDES)

faqs = [
    ("Q", "Saya tidak bisa login. Apa yang harus dilakukan?",
     "Periksa username (ketua.xxx) dan password. Pastikan Caps Lock nonaktif. Jika tetap gagal, hubungi Admin Sistem untuk reset password atau aktivasi akun."),
    ("Q", "Mengapa tombol Submit tidak muncul/tidak aktif?",
     "Tombol Submit hanya aktif jika SEMUA IKU sudah diisi target TW I–IV dan minimal satu rencana kegiatan. Periksa apakah ada IKU yang belum diisi."),
    ("Q", "Bisakah saya mengedit data setelah submit?",
     "Tidak. Dokumen yang telah disubmit dikunci secara otomatis. Pengeditan hanya bisa dilakukan jika Kabag Umum mengembalikan dokumen dengan catatan revisi."),
    ("Q", "Apa itu IKU Kolaborasi? Siapa yang mengisinya?",
     "IKU Kolaborasi adalah IKU yang dikerjakan bersama beberapa Tim Kerja. Setiap Tim Kerja mengisi rencana kegiatannya masing-masing pada IKU tersebut."),
    ("Q", "Berapa lama proses persetujuan oleh Kabag Umum?",
     "Proses persetujuan dilakukan oleh Kabag Umum sesuai jadwal yang ditetapkan. Pantau status dokumen secara berkala di halaman Penyusunan."),
    ("Q", "Data yang saya input hilang/tidak tersimpan. Apa yang terjadi?",
     "Pastikan Anda selalu menekan tombol Simpan setelah mengisi data. Data tidak tersimpan otomatis. Jika masalah berlanjut, hubungi Admin Sistem."),
]

cols2 = 2
rows2 = 3
for i, (q_label, question, answer) in enumerate(faqs):
    col = i % cols2
    row = i // cols2
    lx = 0.4 + col * 6.5
    ty = 1.38 + row * 1.88
    add_rect(slide, lx, ty, 6.1, 1.8, fill_color=PUTIH,
             line_color=RGBColor(0xCC, 0xD9, 0xEA), line_width=Pt(1))
    add_rect(slide, lx, ty, 0.55, 1.8, fill_color=BIRU_MUDA)
    add_text_box(slide, q_label, lx + 0.07, ty + 0.65, 0.42, 0.55,
                 font_size=Pt(18), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    add_text_box(slide, question, lx + 0.65, ty + 0.05, 5.35, 0.48,
                 font_size=Pt(10), bold=True, color=BIRU_TUA, wrap=True)
    add_rect(slide, lx + 0.65, ty + 0.54, 5.35, 0.02, fill_color=RGBColor(0xCC, 0xD9, 0xEA))
    add_text_box(slide, answer, lx + 0.65, ty + 0.62, 5.35, 1.1,
                 font_size=Pt(9.5), color=ABU_TUA, wrap=True)

# Penutup strip
add_rect(slide, 0.4, 7.05, 12.5, 0.45, fill_color=BIRU_TUA)
add_text_box(slide, "Untuk bantuan lebih lanjut, hubungi Admin Sistem SIPITUNG  |  LLDIKTI Wilayah III  |  Versi 1.0 — April 2026",
             0.6, 7.07, 12.2, 0.4, font_size=Pt(9.5), color=PUTIH, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
output_path = "Juknis_SIPITUNG_Ketua_Tim_Kerja.pptx"
prs.save(output_path)
print(f"✓ File berhasil dibuat: {output_path}")
print(f"  Total slides: {prs.slides.__len__()}")
