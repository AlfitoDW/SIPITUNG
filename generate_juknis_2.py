"""
Script untuk generate Juknis SIPITUNG tambahan:
  - Juknis_SIPITUNG_Ketua_Tim_Kerja_Pengukuran.pptx
  - Juknis_SIPITUNG_Kabag_Umum.pptx

Jalankan: python3 generate_juknis_2.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ─── Palet warna ────────────────────────────────────────────────────────────
BIRU_TUA   = RGBColor(0x1A, 0x37, 0x6C)
BIRU_MID   = RGBColor(0x1E, 0x4D, 0x9B)
BIRU_MUDA  = RGBColor(0x2E, 0x75, 0xB6)
HIJAU_TUA  = RGBColor(0x1B, 0x6B, 0x35)
HIJAU_MID  = RGBColor(0x27, 0xAE, 0x60)
HIJAU_MUDA = RGBColor(0xD5, 0xF5, 0xE3)
AKSEN      = RGBColor(0xF5, 0xA6, 0x23)
MERAH      = RGBColor(0xC0, 0x39, 0x2B)
AMBER_BG   = RGBColor(0xFF, 0xF3, 0xCD)
AMBER_BD   = RGBColor(0xF5, 0xA6, 0x23)
MERAH_BG   = RGBColor(0xFD, 0xED, 0xEC)
HIJAU_BG   = RGBColor(0xE8, 0xF5, 0xE9)
PUTIH      = RGBColor(0xFF, 0xFF, 0xFF)
ABU_MUDA   = RGBColor(0xF2, 0xF4, 0xF8)
ABU_TUA    = RGBColor(0x44, 0x44, 0x44)
PLACEHOLDER_BG = RGBColor(0xE8, 0xEE, 0xF8)

BLANK = None  # assigned per prs

# ─── Helper functions ────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.width = lw
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h, sz=Pt(12), bold=False,
       color=ABU_TUA, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = sz; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return box

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.1, fill=BIRU_TUA)
    rect(slide, 0, 1.1, 13.33, 0.06, fill=AKSEN)
    tb(slide, title, 0.35, 0.1, 12.5, 0.68, sz=Pt(24), bold=True, color=PUTIH)
    if subtitle:
        tb(slide, subtitle, 0.35, 0.7, 12.5, 0.38, sz=Pt(12),
           color=RGBColor(0xC8, 0xD8, 0xF0))

def footer(slide, n, total, label="SIPITUNG — Petunjuk Teknis | LLDIKTI Wilayah III"):
    rect(slide, 0, 7.1, 13.33, 0.4, fill=BIRU_TUA)
    tb(slide, label, 0.35, 7.1, 10, 0.4, sz=Pt(9),
       color=RGBColor(0xC8, 0xD8, 0xF0))
    tb(slide, f"{n} / {total}", 12.5, 7.1, 0.8, 0.4, sz=Pt(9),
       color=PUTIH, align=PP_ALIGN.RIGHT)

def step_row(slide, num, title, desc, lx, ty, w=6.6, h=0.85, accent=BIRU_MID):
    rect(slide, lx, ty, w, h, fill=PUTIH,
         line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(0.75))
    rect(slide, lx, ty, 0.5, h, fill=accent)
    tb(slide, str(num), lx, ty + h/2 - 0.22, 0.5, 0.44,
       sz=Pt(17), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, title, lx + 0.58, ty + 0.05, w - 0.65, 0.3,
       sz=Pt(10.5), bold=True, color=BIRU_TUA)
    tb(slide, desc,  lx + 0.58, ty + 0.36, w - 0.65, h - 0.42,
       sz=Pt(9.5), color=ABU_TUA, wrap=True)

def ph(slide, caption, lx, ty, w, h):
    """Placeholder screenshot box."""
    rect(slide, lx, ty, w, h, fill=PLACEHOLDER_BG,
         line=BIRU_MUDA, lw=Pt(1.5))
    rect(slide, lx, ty, w, 0.38, fill=BIRU_MUDA)
    tb(slide, caption, lx + 0.15, ty + 0.02, w - 0.2, 0.35,
       sz=Pt(10), bold=True, color=PUTIH)
    inner = caption.replace("Tampilan ", "").replace("Tombol ", "")
    tb(slide, f"[ SCREENSHOT {inner.upper()} ]\n\nTambahkan screenshot di sini",
       lx + 0.3, ty + 0.55, w - 0.45, h - 0.7,
       sz=Pt(10), color=BIRU_MUDA, align=PP_ALIGN.CENTER, italic=True)

def note(slide, text, lx, ty, w, h, bg=AMBER_BG, bd=AMBER_BD, fc=RGBColor(0x7D, 0x60, 0x08)):
    rect(slide, lx, ty, w, h, fill=bg, line=bd, lw=Pt(1))
    tb(slide, text, lx + 0.15, ty + 0.04, w - 0.22, h - 0.08,
       sz=Pt(9), color=fc, wrap=True)

def toc_slide(prs, items, title_text, pg, total, col1_color=BIRU_MID):
    slide = blank_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
    header(slide, title_text)
    footer(slide, pg, total)
    for i, (num, ttl, dsc) in enumerate(items):
        col = i % 2; row = i // 2
        lx = 0.4 + col * 6.5; ty = 1.35 + row * 1.12
        rect(slide, lx, ty, 6.1, 1.0, fill=PUTIH,
             line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
        rect(slide, lx, ty, 0.7, 1.0, fill=col1_color)
        tb(slide, num, lx + 0.05, ty + 0.2, 0.6, 0.55,
           sz=Pt(16), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
        tb(slide, ttl, lx + 0.78, ty + 0.06, 5.2, 0.32,
           sz=Pt(11), bold=True, color=BIRU_TUA)
        tb(slide, dsc, lx + 0.78, ty + 0.38, 5.2, 0.55,
           sz=Pt(9), color=ABU_TUA, wrap=True)
    return slide

def cover_slide(prs, nama_pengguna, modul_text, sub_text, accent_color=BIRU_MID):
    slide = blank_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=BIRU_TUA)
    rect(slide, 0, 0, 4.5, 7.5, fill=RGBColor(0x0F, 0x20, 0x44))
    rect(slide, 4.5, 0, 0.08, 7.5, fill=AKSEN)
    # Logo placeholder
    rect(slide, 0.4, 0.4, 3.7, 1.6, fill=RGBColor(0x1A, 0x37, 0x6C),
         line=AKSEN, lw=Pt(1.5))
    tb(slide, "[LOGO LLDIKTI\nWILAYAH III]", 0.4, 0.4, 3.7, 1.6,
       sz=Pt(11), color=RGBColor(0xC8, 0xD8, 0xF0), align=PP_ALIGN.CENTER)
    tb(slide, "LLDIKTI WILAYAH III", 0.3, 2.3, 3.9, 0.55,
       sz=Pt(13), bold=True, color=AKSEN)
    tb(slide, "Kementerian Pendidikan Tinggi,\nSains dan Teknologi",
       0.3, 2.8, 3.9, 0.7, sz=Pt(10), color=RGBColor(0xC8, 0xD8, 0xF0))
    rect(slide, 0.3, 3.65, 3.9, 0.05, fill=AKSEN)
    tb(slide, "PETUNJUK TEKNIS", 0.3, 3.8, 3.9, 0.4,
       sz=Pt(11), bold=True, color=RGBColor(0xC8, 0xD8, 0xF0))
    tb(slide, "Penggunaan Aplikasi\nSIPITUNG",
       0.3, 4.2, 3.9, 0.9, sz=Pt(14), bold=True, color=PUTIH)
    tb(slide, "Versi 1.0  |  April 2026", 0.3, 6.9, 3.9, 0.4,
       sz=Pt(9), color=RGBColor(0x88, 0xA8, 0xD0))
    # Kanan
    tb(slide, "SIPITUNG", 4.9, 1.1, 8.0, 1.4,
       sz=Pt(52), bold=True, color=PUTIH)
    tb(slide, "Sistem Informasi Perencanaan & Kinerja",
       4.9, 2.45, 8.0, 0.5, sz=Pt(16), color=AKSEN)
    rect(slide, 4.9, 3.05, 8.0, 0.05, fill=RGBColor(0x88, 0xA8, 0xD0))
    tb(slide, modul_text, 4.9, 3.2, 8.0, 0.5,
       sz=Pt(15), bold=True, color=PUTIH)
    tb(slide, sub_text, 4.9, 3.75, 8.0, 0.8,
       sz=Pt(13), color=RGBColor(0xC8, 0xD8, 0xF0))
    chip_data = [("Pengguna", nama_pengguna), ("Tahun", "2026"), ("Status", "Aktif")]
    for i, (lbl, val) in enumerate(chip_data):
        cx = 4.9 + i * 2.7
        rect(slide, cx, 4.7, 2.4, 0.65, fill=accent_color,
             line=RGBColor(0x88, 0xA8, 0xD0), lw=Pt(0.75))
        tb(slide, lbl, cx + 0.1, 4.72, 2.2, 0.25, sz=Pt(8), bold=True, color=AKSEN)
        tb(slide, val, cx + 0.1, 4.97, 2.2, 0.3,  sz=Pt(11), bold=True, color=PUTIH)
    return slide


# ════════════════════════════════════════════════════════════════════════════
# FILE 1 — Juknis Ketua Tim Kerja: Pengukuran Kinerja
# ════════════════════════════════════════════════════════════════════════════
prs1 = new_prs()
T1 = 12  # total slides

# ── Slide 1: Cover ──────────────────────────────────────────────────────────
cover_slide(prs1, "Ketua Tim Kerja", "Modul Pengukuran Kinerja",
            "Panduan Pengisian Realisasi & Pelaporan", accent_color=HIJAU_TUA)

# ── Slide 2: Daftar Isi ─────────────────────────────────────────────────────
toc_items_1 = [
    ("01", "Pendahuluan & Tujuan",       "Latar belakang modul Pengukuran Kinerja"),
    ("02", "Ketentuan Umum",             "Aturan penting sebelum mengisi realisasi"),
    ("03", "Login & Navigasi",           "Masuk ke menu Pengukuran Kinerja"),
    ("04", "Memahami Halaman Utama",     "Struktur tabel IKU dan kolom realisasi"),
    ("05", "Mengisi Data Realisasi",     "Cara mengisi nilai realisasi per IKU"),
    ("06", "IKU Kolaborasi (Co-PIC)",    "Koordinasi pengisian bersama tim lain"),
    ("07", "Submit Laporan ke Kabag",    "Cara mengirim laporan triwulan"),
    ("08", "Revisi & Submit Ulang",      "Prosedur jika laporan dikembalikan"),
    ("09", "Status & Notifikasi",        "Memantau status persetujuan laporan"),
    ("10", "FAQ",                        "Pertanyaan yang sering diajukan"),
]
toc_slide(prs1, toc_items_1, "Daftar Isi — Juknis Pengukuran Kinerja", 2, T1, col1_color=HIJAU_MID)

# ── Slide 3: Pendahuluan & Ketentuan ────────────────────────────────────────
slide = blank_slide(prs1)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "01  |  Pendahuluan & Tujuan", "Konteks dan maksud modul Pengukuran Kinerja")
footer(slide, 3, T1)

tb(slide, "Tentang Modul Pengukuran Kinerja",
   0.4, 1.38, 6.0, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 0.4, 1.38, 0.06, 0.38, fill=HIJAU_MID)
tb(slide,
   "Modul Pengukuran Kinerja memungkinkan setiap Ketua Tim Kerja untuk mengisi "
   "data realisasi capaian Indikator Kinerja Utama (IKU) per triwulan. "
   "Data ini menjadi dasar evaluasi kinerja dan pelaporan kepada Kabag Umum. "
   "Pengisian dilakukan setelah Admin Sistem mengaktifkan periode triwulan yang bersangkutan.",
   0.5, 1.85, 5.9, 1.5, sz=Pt(10.5), color=ABU_TUA, wrap=True)

tb(slide, "Tujuan", 0.4, 3.5, 6.0, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 0.4, 3.5, 0.06, 0.38, fill=HIJAU_MID)
goals = [
    "Memandu Ketua Tim Kerja dalam mengisi data realisasi IKU per triwulan.",
    "Memastikan data capaian kinerja tercatat secara akurat dan tepat waktu.",
    "Menjelaskan alur koordinasi untuk IKU yang dikerjakan bersama (Co-PIC).",
    "Menjadi panduan dalam proses submit laporan hingga mendapat persetujuan.",
]
for i, g in enumerate(goals):
    ty = 4.0 + i * 0.54
    rect(slide, 0.5, ty + 0.07, 0.24, 0.24, fill=HIJAU_MID)
    tb(slide, str(i+1), 0.5, ty + 0.04, 0.24, 0.28,
       sz=Pt(9), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, g, 0.82, ty, 5.6, 0.5, sz=Pt(10.5), color=ABU_TUA, wrap=True)

# Panel kanan — siklus triwulan
rect(slide, 7.0, 1.35, 6.0, 5.7, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 7.0, 1.35, 6.0, 0.45, fill=HIJAU_TUA)
tb(slide, "Siklus Pengukuran Kinerja Per Tahun",
   7.15, 1.37, 5.7, 0.4, sz=Pt(11), bold=True, color=PUTIH)

tw_data = [
    ("TW I",  "Januari – Maret",    "Laporan capaian Q1"),
    ("TW II", "April – Juni",       "Laporan capaian Q2"),
    ("TW III","Juli – September",   "Laporan capaian Q3"),
    ("TW IV", "Oktober – Desember", "Laporan capaian Q4 + evaluasi akhir tahun"),
]
for i, (tw, bulan, ket) in enumerate(tw_data):
    ty = 1.95 + i * 1.22
    bg = HIJAU_MUDA if i % 2 == 0 else PUTIH
    rect(slide, 7.0, ty, 6.0, 1.18, fill=bg)
    rect(slide, 7.0, ty, 1.0, 1.18, fill=HIJAU_MID if i % 2 == 0 else HIJAU_TUA)
    tb(slide, tw, 7.0, ty + 0.35, 1.0, 0.5,
       sz=Pt(13), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, bulan, 8.1, ty + 0.1, 4.7, 0.35,
       sz=Pt(11), bold=True, color=BIRU_TUA)
    tb(slide, ket, 8.1, ty + 0.48, 4.7, 0.6,
       sz=Pt(9.5), color=ABU_TUA, wrap=True)

# ── Slide 4: Ketentuan Umum ──────────────────────────────────────────────────
slide = blank_slide(prs1)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "02  |  Ketentuan Umum", "Aturan yang wajib dipahami sebelum mengisi data")
footer(slide, 4, T1)

ku = [
    ("1", "Periode Aktif", "Pengisian realisasi hanya bisa dilakukan saat Admin Sistem telah mengaktifkan periode triwulan yang bersangkutan. Jika tabel IKU tidak tampil, kemungkinan periode belum diaktifkan."),
    ("2", "IKU yang Tampil", "Hanya IKU di mana Tim Kerja Anda tercatat sebagai PIC (primary atau co-PIC) yang akan tampil di halaman Pengukuran Kinerja."),
    ("3", "IKU Kolaborasi", "Untuk IKU yang dimiliki bersama tim lain (co-PIC), semua PIC dapat mengisi dan mengedit realisasi. Data terakhir yang disimpan akan menjadi nilai yang tersimpan (last-save-wins)."),
    ("4", "Submit Laporan", "Setelah semua IKU terisi, ajukan laporan ke Kabag Umum. Hanya SATU tim yang dapat submit untuk IKU kolaborasi — tim lain tidak perlu submit ulang untuk IKU yang sama."),
    ("5", "Laporan Tersubmit", "Laporan yang sudah disubmit tidak dapat diubah kecuali dikembalikan oleh Kabag Umum."),
    ("6", "Kelengkapan Isian", "Untuk setiap IKU, setidaknya nilai realisasi wajib diisi. Kolom progress kegiatan, kendala, dan strategi tindak lanjut sangat dianjurkan untuk diisi secara lengkap."),
]
cols = 2
for i, (num, ttl, dsc) in enumerate(ku):
    col = i % cols; row = i // cols
    lx = 0.4 + col * 6.5; ty = 1.38 + row * 1.92
    rect(slide, lx, ty, 6.1, 1.85, fill=PUTIH,
         line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
    rect(slide, lx, ty, 6.1, 0.08, fill=HIJAU_MID)
    rect(slide, lx, ty + 0.08, 0.55, 1.77, fill=HIJAU_MUDA)
    tb(slide, num, lx + 0.07, ty + 0.6, 0.4, 0.55,
       sz=Pt(22), bold=True, color=HIJAU_TUA, align=PP_ALIGN.CENTER)
    tb(slide, ttl, lx + 0.65, ty + 0.12, 5.3, 0.32,
       sz=Pt(11), bold=True, color=BIRU_TUA)
    rect(slide, lx + 0.65, ty + 0.46, 5.35, 0.02, fill=RGBColor(0xCC, 0xD9, 0xEA))
    tb(slide, dsc, lx + 0.65, ty + 0.52, 5.35, 1.26,
       sz=Pt(9.5), color=ABU_TUA, wrap=True)

# ── Slide 5: Login & Navigasi ────────────────────────────────────────────────
slide = blank_slide(prs1)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "03  |  Login & Navigasi ke Pengukuran Kinerja")
footer(slide, 5, T1)

# Breadcrumb
rect(slide, 0.4, 1.35, 12.5, 0.52, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
bcs = ["Dashboard", "Pengukuran Kinerja"]
for i, bc in enumerate(bcs):
    bx = 0.5 + i * 3.1
    bg = HIJAU_TUA if i == len(bcs)-1 else RGBColor(0xE0, 0xE8, 0xF5)
    fc = PUTIH if i == len(bcs)-1 else BIRU_TUA
    rect(slide, bx, 1.38, 2.8, 0.46, fill=bg)
    tb(slide, ("▶  " if i > 0 else "") + bc, bx + 0.1, 1.41, 2.6, 0.4,
       sz=Pt(10), bold=(i == len(bcs)-1), color=fc)

# Screenshot placeholder
ph(slide, "Tampilan Menu Sidebar Pengukuran", 0.4, 2.0, 4.5, 5.0)

# Langkah
tb(slide, "Langkah Login & Navigasi",
   5.25, 2.0, 7.7, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.25, 2.0, 0.06, 0.38, fill=AKSEN)

nav_steps = [
    ("1", "Buka Aplikasi & Login",
     'Buka browser, akses URL SIPITUNG. Masukkan username (contoh: ketua.pk) dan password, lalu klik "Masuk".'),
    ("2", "Pilih Menu Pengukuran Kinerja",
     "Pada sidebar kiri, klik menu 'Pengukuran Kinerja' (ikon grafik batang). Menu ini langsung mengarah ke halaman pengisian."),
    ("3", "Pilih Periode Aktif",
     "Di halaman Pengukuran, pilih periode triwulan yang aktif dari dropdown periode (contoh: Triwulan I 2026)."),
    ("4", "Verifikasi Daftar IKU",
     "Pastikan daftar IKU Tim Kerja Anda tampil dengan benar. Jika kosong, kemungkinan periode belum diaktifkan atau belum ada IKU yang di-assign ke Tim Anda."),
    ("5", "Mulai Isi Realisasi",
     "Tabel siap diisi. Lanjutkan ke langkah berikutnya untuk cara mengisi nilai realisasi."),
]
for i, (num, ttl, dsc) in enumerate(nav_steps):
    step_row(slide, num, ttl, dsc, 5.25, 2.5 + i * 0.92, w=7.7, h=0.85, accent=HIJAU_MID)

# ── Slide 6: Memahami Halaman Utama ─────────────────────────────────────────
slide = blank_slide(prs1)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "04  |  Memahami Halaman Pengukuran Kinerja",
       "Penjelasan setiap kolom dan tombol yang tersedia")
footer(slide, 6, T1)

ph(slide, "Tampilan Halaman Utama Pengukuran", 0.4, 1.35, 5.2, 5.75)

# Penjelasan kolom
tb(slide, "Penjelasan Kolom Tabel",
   5.95, 1.38, 7.1, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.95, 1.38, 0.06, 0.38, fill=AKSEN)

cols_info = [
    ("Sasaran",       "Nama sasaran strategis yang menjadi kelompok IKU (color-coded per sasaran)."),
    ("Indikator",     "Nama dan kode IKU. Ikon 👥 muncul jika IKU dimiliki bersama tim lain (Co-PIC)."),
    ("PIC",           "Tim Kerja penanggung jawab utama IKU. Badge biru = primer, abu = co-PIC."),
    ("Satuan",        "Satuan pengukuran IKU (%, orang, dokumen, dsb.)."),
    ("Target PK",     "Target tahunan sesuai Perjanjian Kinerja."),
    ("Target TW",     "Target triwulan sesuai Rencana Aksi yang telah disetujui."),
    ("Realisasi",     "Nilai capaian aktual yang Anda isi. Klik tombol 'Isi' atau 'Edit' untuk mengisi."),
    ("Status",        "Menunjukkan apakah realisasi sudah diisi (✓ Terisi) atau belum (○ Belum)."),
    ("Aksi",          "Tombol 'Isi' (data belum ada) atau 'Edit' (data sudah ada) untuk membuka form isian."),
]
for i, (col, dsc) in enumerate(cols_info):
    ty = 1.9 + i * 0.59
    bg = ABU_MUDA if i % 2 == 0 else PUTIH
    rect(slide, 5.95, ty, 7.1, 0.57, fill=bg)
    tb(slide, col, 6.1, ty + 0.06, 1.6, 0.42,
       sz=Pt(10), bold=True, color=BIRU_TUA)
    rect(slide, 7.68, ty + 0.1, 0.03, 0.38, fill=RGBColor(0xCC, 0xD9, 0xEA))
    tb(slide, dsc, 7.78, ty + 0.06, 5.2, 0.48,
       sz=Pt(9.5), color=ABU_TUA, wrap=True)

# ── Slide 7: Mengisi Data Realisasi ─────────────────────────────────────────
slide = blank_slide(prs1)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "05  |  Mengisi Data Realisasi IKU",
       "Klik tombol Isi/Edit pada setiap baris IKU untuk membuka form isian")
footer(slide, 7, T1)

ph(slide, "Tampilan Form Isian Realisasi", 0.4, 1.35, 5.2, 5.75)

tb(slide, "Cara Mengisi Data Realisasi",
   5.95, 1.38, 7.1, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.95, 1.38, 0.06, 0.38, fill=AKSEN)

realisasi_steps = [
    ("1", "Klik Tombol 'Isi' atau 'Edit'",
     "Pada baris IKU yang ingin diisi, klik tombol 'Isi' (data belum ada) atau 'Edit' (data sudah ada). Form dialog akan terbuka."),
    ("2", "Isi Nilai Realisasi (Wajib)",
     "Masukkan nilai capaian aktual untuk periode ini. Gunakan angka sesuai satuan IKU (contoh: 85 untuk 85%)."),
    ("3", "Isi Progress Kegiatan",
     "Tuliskan deskripsi singkat progres kegiatan yang telah dilaksanakan dalam triwulan ini."),
    ("4", "Isi Kendala (jika ada)",
     "Jika ada hambatan dalam pencapaian target, tuliskan secara singkat dan jelas."),
    ("5", "Isi Strategi Tindak Lanjut",
     "Tuliskan rencana tindakan untuk mengatasi kendala atau meningkatkan capaian di triwulan berikutnya."),
    ("6", "Klik Simpan",
     'Klik tombol "Simpan". Data tersimpan dan status IKU berubah menjadi ✓ Terisi. Ulangi untuk semua IKU.'),
]
for i, (num, ttl, dsc) in enumerate(realisasi_steps):
    step_row(slide, num, ttl, dsc, 5.95, 1.9 + i * 0.85, w=7.1, h=0.78, accent=HIJAU_MID)

note(slide, "✓  Tips: Pastikan nilai realisasi sesuai kondisi nyata. Data ini akan direview oleh Kabag Umum dan menjadi dasar evaluasi kinerja Tim Kerja Anda.",
     5.95, 6.95, 7.1, 0.55, bg=HIJAU_BG, bd=RGBColor(0x27, 0xAE, 0x60), fc=HIJAU_TUA)

# ── Slide 8: IKU Kolaborasi ──────────────────────────────────────────────────
slide = blank_slide(prs1)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "06  |  IKU Kolaborasi (Co-PIC)",
       "Pengisian IKU yang dikerjakan bersama tim kerja lain")
footer(slide, 8, T1)

# Panel kiri — penjelasan
rect(slide, 0.4, 1.35, 5.8, 5.75, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 0.4, 1.35, 5.8, 0.45, fill=BIRU_MID)
tb(slide, "Apa itu IKU Kolaborasi?",
   0.55, 1.37, 5.5, 0.4, sz=Pt(12), bold=True, color=PUTIH)

tb(slide,
   "IKU Kolaborasi adalah Indikator Kinerja Utama yang menjadi tanggung jawab "
   "lebih dari satu Tim Kerja. Ikon 👥 ditampilkan pada nama IKU tersebut, "
   "dan banner koordinasi akan muncul di atas tabel.\n\n"
   "Semua PIC dapat mengisi dan mengedit nilai realisasi kapan saja. "
   "Data yang terakhir disimpan akan menjadi nilai final (last-save-wins).",
   0.55, 1.88, 5.6, 1.55, sz=Pt(10.5), color=ABU_TUA, wrap=True)

aturan_kolab = [
    ("Siapa yang bisa isi?", "Semua Tim Kerja yang tercatat sebagai PIC (primer atau co-PIC) dapat mengisi dan mengedit realisasi IKU kolaborasi."),
    ("Last-save-wins?", "Nilai terakhir yang disimpan akan menjadi nilai final. Koordinasikan dengan tim lain sebelum menyimpan jika ada perbedaan data."),
    ("Catatan koordinasi?", "Gunakan kolom 'Catatan' (berlabel 'terlihat oleh semua PIC') untuk komunikasi internal antar tim."),
    ("Submit laporan?", "Hanya SATU tim yang perlu submit laporan untuk IKU kolaborasi. Tim lain akan mendapat notifikasi bahwa kolaborator sudah submit."),
]
for i, (q, a) in enumerate(aturan_kolab):
    ty = 3.55 + i * 0.84
    rect(slide, 0.4, ty, 5.8, 0.8, fill=ABU_MUDA if i % 2 == 0 else PUTIH)
    tb(slide, q, 0.55, ty + 0.05, 5.6, 0.28, sz=Pt(10), bold=True, color=BIRU_TUA)
    tb(slide, a, 0.55, ty + 0.35, 5.6, 0.42, sz=Pt(9.5), color=ABU_TUA, wrap=True)

# Panel kanan — banner koordinasi
rect(slide, 6.6, 1.35, 6.35, 5.75, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 6.6, 1.35, 6.35, 0.45, fill=BIRU_TUA)
tb(slide, "Elemen UI untuk IKU Kolaborasi",
   6.75, 1.37, 6.1, 0.4, sz=Pt(12), bold=True, color=PUTIH)

ui_elems = [
    ("Banner Koordinasi\n(di atas tabel)",
     "Muncul otomatis jika ada IKU shared. Menampilkan daftar IKU yang dikerjakan bersama dan menyoroti 'Tim Anda' dengan badge."),
    ("Ikon 👥 pada Nama IKU",
     "Penanda visual bahwa IKU ini dikerjakan bersama. Badge nama singkat semua PIC ditampilkan di form dialog."),
    ("Kolom 'Terisi Oleh'",
     "Menampilkan nama singkat Tim Kerja yang terakhir menyimpan data realisasi untuk IKU tersebut."),
    ("Catatan Koordinasi\n(di dalam form)",
     "Kolom opsional dalam form isian. Catatan ini terlihat oleh semua PIC dan berguna untuk koordinasi internal."),
    ("Banner Info Submit\n(jika kolaborator sudah submit)",
     "Jika tim lain sudah submit laporan terlebih dahulu, sebuah banner amber muncul sebagai informasi. Tim Anda tetap dapat submit secara mandiri."),
]
for i, (ttl, dsc) in enumerate(ui_elems):
    ty = 1.9 + i * 0.98
    rect(slide, 6.6, ty, 6.35, 0.94, fill=ABU_MUDA if i % 2 == 0 else PUTIH)
    tb(slide, ttl, 6.75, ty + 0.05, 2.0, 0.5, sz=Pt(10), bold=True, color=BIRU_MID)
    rect(slide, 8.72, ty + 0.1, 0.03, 0.72, fill=RGBColor(0xCC, 0xD9, 0xEA))
    tb(slide, dsc, 8.82, ty + 0.08, 4.0, 0.8, sz=Pt(9.5), color=ABU_TUA, wrap=True)

note(slide, "⚠  Koordinasikan dengan tim co-PIC sebelum menyimpan data akhir agar tidak terjadi perbedaan nilai yang tidak disengaja.",
     0.4, 7.05, 12.5, 0.45)

# ── Slide 9: Submit Laporan ──────────────────────────────────────────────────
slide = blank_slide(prs1)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "07  |  Submit Laporan Pengukuran ke Kabag Umum",
       "Mengirim laporan triwulan setelah semua IKU terisi")
footer(slide, 9, T1)

# Alur
rect(slide, 0.4, 1.35, 12.5, 1.42, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
alur = [
    ("Isi Semua\nRealisasi IKU", HIJAU_MID, "Ketua Tim"),
    ("Klik Submit\nLaporan",     BIRU_MID,  "Ketua Tim"),
    ("Review\nKabag Umum",       BIRU_TUA,  "Kabag Umum"),
    ("Laporan\nDisetujui",       HIJAU_TUA, "Kabag Umum"),
]
for i, (lbl, col, actor) in enumerate(alur):
    ax = 0.6 + i * 3.1
    rect(slide, ax, 1.42, 2.7, 1.28, fill=col)
    tb(slide, lbl, ax + 0.1, 1.47, 2.5, 0.6,
       sz=Pt(10), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, actor, ax + 0.1, 2.05, 2.5, 0.55,
       sz=Pt(8.5), color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
    if i < len(alur)-1:
        tb(slide, "▶", ax + 2.73, 1.85, 0.35, 0.55,
           sz=Pt(16), color=BIRU_MUDA, align=PP_ALIGN.CENTER)

ph(slide, "Tampilan Tombol Submit Laporan", 0.4, 2.9, 4.5, 4.15)

tb(slide, "Cara Submit Laporan",
   5.25, 2.9, 7.7, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.25, 2.9, 0.06, 0.38, fill=AKSEN)

submit_steps = [
    ("1", "Periksa Progress Bar",
     "Pastikan semua langkah di progress bar sudah terpenuhi: Periode aktif ✓ | IKU tersedia ✓ | Realisasi diisi ✓."),
    ("2", "Klik Tombol Submit Laporan",
     'Klik tombol "Submit Laporan ke Kabag Umum" yang tersedia di bawah tabel atau di bagian atas halaman.'),
    ("3", "Konfirmasi Submit",
     "Dialog konfirmasi akan muncul. Baca dengan seksama, pastikan Anda yakin data sudah lengkap dan benar, lalu klik Konfirmasi."),
    ("4", "Pantau Status",
     "Setelah submit berhasil, status laporan berubah menjadi 'Menunggu persetujuan Kabag Umum'. Pantau statusnya secara berkala."),
]
for i, (num, ttl, dsc) in enumerate(submit_steps):
    step_row(slide, num, ttl, dsc, 5.25, 3.4 + i * 0.85, w=7.7, h=0.78, accent=HIJAU_TUA)

note(slide, "⛔  Penting: Laporan yang telah disubmit TIDAK DAPAT diedit. Pastikan semua data sudah benar sebelum submit.",
     5.25, 6.88, 7.7, 0.55, bg=MERAH_BG, bd=MERAH, fc=MERAH)

# ── Slide 10: Revisi & Submit Ulang ─────────────────────────────────────────
slide = blank_slide(prs1)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "08–09  |  Revisi, Status & Notifikasi",
       "Prosedur jika laporan dikembalikan dan cara memantau status")
footer(slide, 10, T1)

# Kiri — Revisi
rect(slide, 0.4, 1.35, 5.9, 5.75, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 0.4, 1.35, 5.9, 0.45, fill=MERAH)
tb(slide, "08  |  Revisi & Submit Ulang",
   0.55, 1.37, 5.6, 0.4, sz=Pt(12), bold=True, color=PUTIH)

revisi_steps = [
    ("Status 'Laporan Dikembalikan'",
     "Jika Kabag Umum mengembalikan laporan, status berubah dan banner merah muncul di halaman Pengukuran Kinerja."),
    ("Baca Catatan Revisi",
     "Klik area catatan/rekomendasi Kabag Umum yang tampil di banner merah untuk melihat poin-poin yang perlu diperbaiki."),
    ("Perbaiki Data Realisasi",
     "Klik tombol 'Edit' pada IKU yang perlu diperbaiki. Ubah nilai atau tambahkan keterangan sesuai catatan Kabag Umum."),
    ("Klik Simpan",
     "Simpan perubahan pada setiap IKU yang diperbaiki."),
    ("Submit Ulang Laporan",
     "Setelah semua perbaikan selesai, klik kembali tombol Submit Laporan. Catatan revisi sebelumnya akan direset otomatis."),
]
for i, (ttl, dsc) in enumerate(revisi_steps):
    ty = 1.9 + i * 0.98
    rect(slide, 0.4, ty, 5.9, 0.94, fill=ABU_MUDA if i % 2 == 0 else PUTIH)
    rect(slide, 0.4, ty, 0.42, 0.94, fill=MERAH)
    tb(slide, str(i+1), 0.4, ty + 0.25, 0.42, 0.44,
       sz=Pt(16), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, ttl, 0.9, ty + 0.06, 5.25, 0.28, sz=Pt(10.5), bold=True, color=BIRU_TUA)
    tb(slide, dsc, 0.9, ty + 0.38, 5.25, 0.52, sz=Pt(9.5), color=ABU_TUA, wrap=True)

# Kanan — Status
rect(slide, 6.7, 1.35, 6.25, 5.75, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 6.7, 1.35, 6.25, 0.45, fill=HIJAU_TUA)
tb(slide, "09  |  Status Laporan & Tampilan Banner",
   6.85, 1.37, 6.0, 0.4, sz=Pt(12), bold=True, color=PUTIH)

statuses = [
    ("Belum Submit",   RGBColor(0xF5, 0xA6, 0x23), "Tombol submit tersedia. Isi semua realisasi terlebih dahulu."),
    ("Menunggu Persetujuan\n(submitted)", RGBColor(0x2E, 0x75, 0xB6),
     "Banner abu muncul: 'Menunggu persetujuan Kabag Umum'. Tidak dapat edit."),
    ("Disetujui\n(kabag_approved)", HIJAU_MID,
     "Banner hijau muncul: 'Laporan Disetujui'. Catatan apresiasi Kabag ditampilkan."),
    ("Dikembalikan\n(rejected)", MERAH,
     "Banner merah muncul: 'Laporan Dikembalikan'. Catatan revisi Kabag ditampilkan. Tombol submit ulang muncul."),
]
for i, (status, col, dsc) in enumerate(statuses):
    ty = 1.9 + i * 1.18
    rect(slide, 6.7, ty, 6.25, 1.14, fill=ABU_MUDA if i % 2 == 0 else PUTIH)
    rect(slide, 6.7, ty, 0.35, 1.14, fill=col)
    tb(slide, status, 7.12, ty + 0.08, 5.7, 0.44, sz=Pt(10.5), bold=True, color=BIRU_TUA)
    rect(slide, 7.12, ty + 0.54, 5.65, 0.02, fill=RGBColor(0xCC, 0xD9, 0xEA))
    tb(slide, dsc, 7.12, ty + 0.6, 5.7, 0.5, sz=Pt(9.5), color=ABU_TUA, wrap=True)

# ── Slide 11: Progress Bar ───────────────────────────────────────────────────
slide = blank_slide(prs1)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "Progress Bar — 5 Langkah Pengukuran Kinerja",
       "Indikator visual kelengkapan laporan triwulan")
footer(slide, 11, T1)

ph(slide, "Tampilan Progress Bar 5 Langkah", 0.4, 1.35, 5.0, 5.75)

tb(slide, "Penjelasan Setiap Langkah Progress",
   5.8, 1.38, 7.15, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.8, 1.38, 0.06, 0.38, fill=AKSEN)

progress_steps = [
    ("✓", "Periode Aktif",
     "Admin Sistem telah mengaktifkan periode triwulan. Jika step ini belum ✓, hubungi Admin Sistem.",
     HIJAU_MID),
    ("✓", "IKU Tersedia",
     "Ada IKU yang di-assign ke Tim Kerja Anda untuk periode ini. Jika belum ✓, hubungi Admin Sistem.",
     HIJAU_MID),
    ("○", "Realisasi Diisi",
     "Semua IKU yang tampil sudah memiliki nilai realisasi. Step ini ✓ setelah semua baris terisi.",
     BIRU_MUDA),
    ("○", "Laporan Disubmit",
     "Laporan sudah dikirim ke Kabag Umum. Step ini ✓ setelah Anda (atau kolaborator co-PIC) submit.",
     BIRU_MID),
    ("○", "Disetujui Kabag",
     "Kabag Umum telah menyetujui laporan. Step ini ✓ setelah status berubah menjadi 'kabag_approved'.",
     HIJAU_TUA),
]
for i, (ico, ttl, dsc, col) in enumerate(progress_steps):
    ty = 1.9 + i * 1.03
    rect(slide, 5.8, ty, 7.15, 0.96, fill=PUTIH,
         line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(0.75))
    rect(slide, 5.8, ty, 0.6, 0.96, fill=col)
    tb(slide, ico, 5.8, ty + 0.26, 0.6, 0.44,
       sz=Pt(18), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, f"Langkah {i+1}: {ttl}", 6.48, ty + 0.06, 6.4, 0.3,
       sz=Pt(10.5), bold=True, color=BIRU_TUA)
    tb(slide, dsc, 6.48, ty + 0.38, 6.4, 0.54,
       sz=Pt(9.5), color=ABU_TUA, wrap=True)

# ── Slide 12: FAQ ────────────────────────────────────────────────────────────
slide = blank_slide(prs1)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "10  |  FAQ — Pengukuran Kinerja", "Pertanyaan yang Sering Diajukan")
footer(slide, 12, T1)

faqs = [
    ("IKU saya tidak tampil di halaman Pengukuran.",
     "Kemungkinan: (1) Periode triwulan belum diaktifkan oleh Admin. (2) Tidak ada IKU yang di-assign ke Tim Anda. Hubungi Admin Sistem."),
    ("Saya sudah mengisi tapi data hilang setelah refresh.",
     "Pastikan Anda sudah klik tombol 'Simpan' setelah mengisi setiap IKU. Data tidak tersimpan otomatis."),
    ("Tombol Submit Laporan tidak muncul/tidak aktif.",
     "Pastikan semua IKU sudah terisi (status ✓ Terisi). Periksa progress bar apakah ada step yang belum terpenuhi."),
    ("Tim lain sudah submit lebih dulu untuk IKU kolaborasi. Apakah saya tetap harus submit?",
     "Tidak wajib, karena laporan sudah diwakili kolaborator. Namun Tim Anda tetap dapat submit secara mandiri jika diperlukan."),
    ("Nilai realisasi yang saya masukkan berbeda dengan milik co-PIC.",
     "Koordinasikan dengan tim co-PIC. Nilai yang terakhir disimpan akan menjadi nilai final (last-save-wins). Gunakan kolom Catatan untuk koordinasi."),
    ("Kabag mengembalikan laporan saya. Apakah data realisasi terhapus?",
     "Tidak. Data realisasi tetap tersimpan. Anda hanya perlu memperbaiki sesuai catatan Kabag, lalu submit ulang."),
]
for i, (q, a) in enumerate(faqs):
    col = i % 2; row = i // 2
    lx = 0.4 + col * 6.5; ty = 1.38 + row * 1.92
    rect(slide, lx, ty, 6.1, 1.85, fill=PUTIH,
         line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
    rect(slide, lx, ty, 6.1, 0.08, fill=HIJAU_MID)
    rect(slide, lx, ty + 0.08, 0.52, 1.77, fill=RGBColor(0xD5, 0xF5, 0xE3))
    tb(slide, "Q", lx + 0.1, ty + 0.66, 0.34, 0.44,
       sz=Pt(18), bold=True, color=HIJAU_TUA, align=PP_ALIGN.CENTER)
    tb(slide, q, lx + 0.62, ty + 0.1, 5.35, 0.48,
       sz=Pt(10), bold=True, color=BIRU_TUA, wrap=True)
    rect(slide, lx + 0.62, ty + 0.6, 5.35, 0.02, fill=RGBColor(0xCC, 0xD9, 0xEA))
    tb(slide, a, lx + 0.62, ty + 0.68, 5.35, 1.12,
       sz=Pt(9.5), color=ABU_TUA, wrap=True)

# Simpan file 1
out1 = "Juknis_SIPITUNG_Ketua_Tim_Kerja_Pengukuran.pptx"
prs1.save(out1)
print(f"✓ File 1 selesai: {out1}  ({len(prs1.slides)} slides)")


# ════════════════════════════════════════════════════════════════════════════
# FILE 2 — Juknis Kabag Umum
# ════════════════════════════════════════════════════════════════════════════
prs2 = new_prs()
T2 = 14

HIJAU_KAB  = RGBColor(0x0D, 0x47, 0x2E)   # warna tema Kabag
HIJAU_KAB2 = RGBColor(0x1A, 0x6B, 0x44)
HIJAU_KAB3 = RGBColor(0x2E, 0x8B, 0x57)

# ── Slide 1: Cover ──────────────────────────────────────────────────────────
cover_slide(prs2, "Kabag Umum (Pimpinan)",
            "Modul Perencanaan & Pengukuran Kinerja",
            "Panduan Review, Persetujuan & Penolakan Dokumen",
            accent_color=HIJAU_KAB)

# ── Slide 2: Daftar Isi ─────────────────────────────────────────────────────
toc_items_2 = [
    ("01", "Pendahuluan & Peran Kabag Umum",  "Tanggung jawab Kabag Umum dalam SIPITUNG"),
    ("02", "Ketentuan Umum",                  "Aturan dan batasan kewenangan"),
    ("03", "Login & Navigasi",                "Cara masuk dan navigasi menu"),
    ("04", "Memahami Tampilan Review",         "Struktur halaman review dokumen"),
    ("05", "Review Rencana Aksi (RA)",         "Cara mereview dan menyetujui/menolak RA"),
    ("06", "Setujui Rencana Aksi",             "Langkah persetujuan RA per Tim Kerja"),
    ("07", "Tolak & Beri Catatan RA",          "Cara menolak RA disertai catatan revisi"),
    ("08", "Review Laporan Pengukuran",        "Cara mereview laporan triwulan Tim Kerja"),
    ("09", "Setujui / Tolak Laporan",          "Langkah persetujuan atau penolakan laporan"),
    ("10", "Export Laporan PDF",               "Cara mengekspor laporan pengukuran per TW"),
    ("11", "Alur Lengkap & Checklist",         "Ringkasan alur kerja dan checklist harian"),
    ("12", "FAQ",                              "Pertanyaan yang sering diajukan"),
]
toc_slide(prs2, toc_items_2, "Daftar Isi — Juknis Kabag Umum", 2, T2, col1_color=HIJAU_KAB2)

# ── Slide 3: Pendahuluan & Peran ─────────────────────────────────────────────
slide = blank_slide(prs2)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "01  |  Pendahuluan & Peran Kabag Umum",
       "Posisi dan tanggung jawab dalam sistem persetujuan SIPITUNG")
footer(slide, 3, T2, "SIPITUNG — Juknis Kabag Umum | LLDIKTI Wilayah III")

# Diagram peran (tengah atas)
rect(slide, 0.4, 1.35, 12.5, 2.5, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 0.4, 1.35, 12.5, 0.42, fill=HIJAU_KAB)
tb(slide, "Posisi Kabag Umum dalam Alur Persetujuan SIPITUNG",
   0.55, 1.37, 12.2, 0.38, sz=Pt(12), bold=True, color=PUTIH)

alur_nodes = [
    ("Tim Kerja\n(Ketua)", BIRU_MID,  "Menyusun RA &\nMengisi Realisasi"),
    ("Kabag Umum\n★ ANDA", HIJAU_KAB, "Review, Setujui/\nTolak Dokumen"),
    ("PPK",                BIRU_TUA,  "Persetujuan\nFinal"),
    ("Terkunci",           HIJAU_MID, "Dokumen\nFinal & Arsip"),
]
for i, (node, col, lbl) in enumerate(alur_nodes):
    nx = 0.8 + i * 3.1
    rect(slide, nx, 1.85, 2.7, 1.85, fill=col)
    tb(slide, node, nx + 0.1, 1.9, 2.5, 0.55,
       sz=Pt(11), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, lbl, nx + 0.1, 2.5, 2.5, 0.62,
       sz=Pt(9.5), color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
    if i < len(alur_nodes)-1:
        tb(slide, "▶", nx + 2.73, 2.5, 0.35, 0.55,
           sz=Pt(16), color=BIRU_MUDA, align=PP_ALIGN.CENTER)

# Bawah: kiri kewenangan, kanan dokumen
rect(slide, 0.4, 4.0, 5.9, 3.1, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 0.4, 4.0, 5.9, 0.42, fill=HIJAU_KAB2)
tb(slide, "Kewenangan Kabag Umum", 0.55, 4.02, 5.65, 0.38,
   sz=Pt(12), bold=True, color=PUTIH)
kewenangan = [
    "Mereview dan menyetujui Rencana Aksi (RA) yang disubmit Tim Kerja",
    "Menolak RA disertai catatan revisi untuk diperbaiki Tim Kerja",
    "Mereview dan menyetujui Laporan Pengukuran Kinerja per triwulan",
    "Menolak Laporan Pengukuran disertai rekomendasi perbaikan",
    "Mengekspor laporan pengukuran dalam format PDF per triwulan",
]
for i, k in enumerate(kewenangan):
    ty = 4.52 + i * 0.5
    rect(slide, 0.4, ty, 0.42, 0.48, fill=HIJAU_KAB3)
    tb(slide, "✓", 0.4, ty + 0.05, 0.42, 0.38,
       sz=Pt(13), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, k, 0.9, ty + 0.05, 5.3, 0.42, sz=Pt(10), color=ABU_TUA, wrap=True)

rect(slide, 6.7, 4.0, 6.25, 3.1, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 6.7, 4.0, 6.25, 0.42, fill=BIRU_TUA)
tb(slide, "Dokumen yang Direview Kabag Umum", 6.85, 4.02, 6.0, 0.38,
   sz=Pt(12), bold=True, color=PUTIH)
docs = [
    ("Modul Perencanaan", "Rencana Aksi (RA) — target triwulan & rencana kegiatan IKU", BIRU_MUDA),
    ("Modul Pengukuran",  "Laporan Pengukuran Kinerja — realisasi capaian per triwulan", HIJAU_MID),
]
for i, (modul, dsc, col) in enumerate(docs):
    ty = 4.52 + i * 1.22
    rect(slide, 6.7, ty, 6.25, 1.15, fill=ABU_MUDA if i % 2 == 0 else PUTIH)
    rect(slide, 6.7, ty, 0.55, 1.15, fill=col)
    tb(slide, modul, 7.34, ty + 0.08, 5.5, 0.32, sz=Pt(10), bold=True, color=BIRU_TUA)
    tb(slide, dsc,   7.34, ty + 0.44, 5.5, 0.6,  sz=Pt(10), color=ABU_TUA, wrap=True)

# ── Slide 4: Ketentuan Umum ──────────────────────────────────────────────────
slide = blank_slide(prs2)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "02  |  Ketentuan Umum", "Aturan dan batasan kewenangan Kabag Umum dalam SIPITUNG")
footer(slide, 4, T2, "SIPITUNG — Juknis Kabag Umum | LLDIKTI Wilayah III")

ku2 = [
    ("1", "Status Dokumen untuk Review",
     "Kabag Umum hanya dapat menyetujui/menolak dokumen berstatus 'submitted'. Dokumen berstatus lain ditampilkan sebagai view-only."),
    ("2", "Persetujuan Bertahap",
     "Setelah Kabag Umum menyetujui, dokumen diteruskan ke PPK untuk persetujuan final. Kabag Umum tidak dapat membatalkan persetujuan PPK."),
    ("3", "Catatan Wajib saat Menolak",
     "Setiap penolakan (RA maupun Laporan) WAJIB disertai catatan/rekomendasi yang jelas agar Tim Kerja dapat melakukan perbaikan."),
    ("4", "View-Only untuk Semua Status",
     "Halaman review menampilkan SEMUA dokumen dari semua Tim Kerja (semua status). Tombol Setujui/Tolak hanya aktif untuk status 'submitted'."),
    ("5", "Login PPK",
     "Perlu diperhatikan: halaman review yang sama digunakan oleh PPK, tetapi tombol aksi hanya muncul untuk akun Kabag Umum (bukan PPK)."),
    ("6", "Export Laporan",
     "Hanya Kabag Umum yang dapat mengekspor laporan pengukuran ke PDF. Export tersedia setelah minimal satu laporan disetujui."),
]
for i, (num, ttl, dsc) in enumerate(ku2):
    col = i % 2; row = i // 2
    lx = 0.4 + col * 6.5; ty = 1.38 + row * 1.92
    rect(slide, lx, ty, 6.1, 1.85, fill=PUTIH,
         line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
    rect(slide, lx, ty, 6.1, 0.08, fill=HIJAU_KAB2)
    rect(slide, lx, ty + 0.08, 0.55, 1.77, fill=RGBColor(0xD5, 0xF5, 0xE3))
    tb(slide, num, lx + 0.1, ty + 0.6, 0.38, 0.55,
       sz=Pt(22), bold=True, color=HIJAU_KAB, align=PP_ALIGN.CENTER)
    tb(slide, ttl, lx + 0.65, ty + 0.12, 5.3, 0.32,
       sz=Pt(11), bold=True, color=BIRU_TUA)
    rect(slide, lx + 0.65, ty + 0.46, 5.35, 0.02, fill=RGBColor(0xCC, 0xD9, 0xEA))
    tb(slide, dsc, lx + 0.65, ty + 0.52, 5.35, 1.26,
       sz=Pt(9.5), color=ABU_TUA, wrap=True)

# ── Slide 5: Login & Navigasi ────────────────────────────────────────────────
slide = blank_slide(prs2)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "03  |  Login & Navigasi Menu Kabag Umum",
       "Cara masuk dan mengakses menu review di SIPITUNG")
footer(slide, 5, T2, "SIPITUNG — Juknis Kabag Umum | LLDIKTI Wilayah III")

ph(slide, "Tampilan Dashboard Kabag Umum", 0.4, 1.35, 4.8, 5.75)

tb(slide, "Langkah Login & Navigasi",
   5.55, 1.38, 7.4, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.55, 1.38, 0.06, 0.38, fill=AKSEN)

login_nav = [
    ("1", "Buka Aplikasi & Login",
     'Buka browser, akses URL SIPITUNG. Masukkan username Kabag Umum (kabag.umum) dan password, lalu klik "Masuk".'),
    ("2", "Dashboard Kabag Umum",
     "Setelah login, Anda akan masuk ke Dashboard yang menampilkan ringkasan: Menunggu Review / Sudah Disetujui / Ditolak."),
    ("3", "Navigasi Modul Perencanaan",
     "Untuk review RA: klik 'Perencanaan' di sidebar → 'Rencana Aksi'. Halaman menampilkan RA dari semua Tim Kerja."),
    ("4", "Navigasi Modul Pengukuran",
     "Untuk review Laporan: klik 'Pengukuran Kinerja' di sidebar. Laporan per Tim Kerja tampil di bagian atas halaman."),
    ("5", "Filter & Pencarian",
     "Gunakan filter tahun anggaran dan pencarian nama tim kerja untuk menemukan dokumen yang ingin direview lebih cepat."),
]
for i, (num, ttl, dsc) in enumerate(login_nav):
    step_row(slide, num, ttl, dsc, 5.55, 1.9 + i * 0.93, w=7.4, h=0.86, accent=HIJAU_KAB2)

# Tabel credentials
rect(slide, 5.55, 6.6, 7.4, 0.9, fill=RGBColor(0xE8, 0xF5, 0xE9),
     line=HIJAU_MID, lw=Pt(1))
tb(slide, "Kredensial Login Kabag Umum:",
   5.7, 6.63, 3.5, 0.28, sz=Pt(10), bold=True, color=HIJAU_KAB)
tb(slide, "Username: kabag.umum\nPassword: (sesuai yang ditetapkan Admin)",
   5.7, 6.9, 7.1, 0.55, sz=Pt(9.5), color=ABU_TUA)

# ── Slide 6: Memahami Tampilan Review ────────────────────────────────────────
slide = blank_slide(prs2)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "04  |  Memahami Tampilan Halaman Review",
       "Struktur tabel dan elemen UI yang digunakan Kabag Umum")
footer(slide, 6, T2, "SIPITUNG — Juknis Kabag Umum | LLDIKTI Wilayah III")

ph(slide, "Tampilan Halaman Review Rencana Aksi", 0.4, 1.35, 5.0, 5.75)

tb(slide, "Penjelasan Kolom Tabel Review",
   5.75, 1.38, 7.2, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.75, 1.38, 0.06, 0.38, fill=AKSEN)

col_info = [
    ("Kolom Status\n(kiri, rowspan tim kerja)", "Nama Tim Kerja (kecil di atas) + badge status berwarna:\nKuning=Draft, Biru=Menunggu Review, Hijau=Disetujui, Merah=Ditolak."),
    ("Kolom Sasaran\n(rowspan per sasaran)", "Nama sasaran strategis, color-coded. Setiap sasaran bisa memiliki beberapa IKU."),
    ("Kolom Indikator", "Kode dan nama IKU. Badge 'Tim Kerja Primer' dan 'Co-PIC' ditampilkan untuk IKU kolaborasi."),
    ("Kolom PIC Tim Kerja", "Badge biru = PIC primer, badge abu = co-PIC. Berguna untuk memastikan IKU kolaborasi sudah diisi kedua pihak."),
    ("Kolom Target / TW I–IV\n(khusus RA)", "Nilai target tahunan dan per triwulan yang diisi Tim Kerja."),
    ("Kolom Aksi\n(hanya untuk kabag_umum)", "Tombol 'Setujui' dan 'Tolak' HANYA muncul saat status = 'submitted'. Status lain menampilkan tanda —."),
]
for i, (k, v) in enumerate(col_info):
    ty = 1.9 + i * 0.88
    bg = ABU_MUDA if i % 2 == 0 else PUTIH
    rect(slide, 5.75, ty, 7.2, 0.84, fill=bg)
    tb(slide, k, 5.9, ty + 0.06, 2.5, 0.5, sz=Pt(10), bold=True, color=BIRU_TUA)
    rect(slide, 8.37, ty + 0.1, 0.03, 0.62, fill=RGBColor(0xCC, 0xD9, 0xEA))
    tb(slide, v, 8.47, ty + 0.06, 4.4, 0.74, sz=Pt(9.5), color=ABU_TUA, wrap=True)

note(slide, "ℹ  Catatan: Halaman ini menampilkan SEMUA Tim Kerja dan SEMUA status. Scroll ke bawah atau gunakan pencarian untuk menemukan dokumen yang perlu ditindaklanjuti.",
     5.75, 7.05, 7.2, 0.45, bg=RGBColor(0xE3, 0xF2, 0xFD), bd=BIRU_MUDA,
     fc=BIRU_TUA)

# ── Slide 7: Review & Setujui RA ─────────────────────────────────────────────
slide = blank_slide(prs2)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "05–06  |  Review & Setujui Rencana Aksi",
       "Cara memeriksa dan memberikan persetujuan RA Tim Kerja")
footer(slide, 7, T2, "SIPITUNG — Juknis Kabag Umum | LLDIKTI Wilayah III")

# Kiri — cara review
rect(slide, 0.4, 1.35, 5.9, 5.75, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 0.4, 1.35, 5.9, 0.45, fill=HIJAU_KAB2)
tb(slide, "05  |  Cara Mereview Rencana Aksi",
   0.55, 1.37, 5.65, 0.4, sz=Pt(12), bold=True, color=PUTIH)

review_points = [
    "Buka halaman Perencanaan → Rencana Aksi. Seluruh RA dari semua Tim Kerja tampil dalam tabel flat.",
    "Identifikasi baris yang memiliki badge status biru 'Menunggu Review' — ini yang perlu ditindaklanjuti.",
    "Periksa kelengkapan target TW I–IV untuk setiap IKU. Pastikan tidak ada nilai yang kosong atau tidak wajar.",
    "Untuk IKU kolaborasi, cek apakah badge PIC primer dan co-PIC sudah mencerminkan Tim Kerja yang benar.",
    "Periksa apakah target triwulan bersifat kumulatif dan sesuai dengan target tahunan PK.",
    "Jika semua data sudah sesuai, klik tombol 'Setujui' pada kolom Aksi baris Tim Kerja tersebut.",
]
for i, pt in enumerate(review_points):
    ty = 1.9 + i * 0.84
    rect(slide, 0.4, ty, 5.9, 0.8, fill=ABU_MUDA if i % 2 == 0 else PUTIH)
    rect(slide, 0.4, ty, 0.42, 0.8, fill=HIJAU_KAB3)
    tb(slide, str(i+1), 0.4, ty + 0.18, 0.42, 0.44,
       sz=Pt(14), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    tb(slide, pt, 0.9, ty + 0.1, 5.28, 0.64, sz=Pt(10), color=ABU_TUA, wrap=True)

# Kanan — cara setujui
rect(slide, 6.7, 1.35, 6.25, 5.75, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 6.7, 1.35, 6.25, 0.45, fill=HIJAU_MID)
tb(slide, "06  |  Langkah Persetujuan RA",
   6.85, 1.37, 6.0, 0.4, sz=Pt(12), bold=True, color=PUTIH)

setujui_steps = [
    ("1", "Temukan Tim Kerja yang Perlu Disetujui",
     "Cari baris dengan status 'Menunggu Review' (badge biru). Tombol Setujui tersedia di kolom Aksi."),
    ("2", "Klik Tombol 'Setujui'",
     "Klik tombol hijau 'Setujui' pada baris Tim Kerja tersebut. Sebuah dialog konfirmasi akan muncul."),
    ("3", "Konfirmasi Persetujuan",
     "Baca isi dialog, lalu klik 'Konfirmasi' untuk menyetujui RA. Status dokumen berubah menjadi 'Menunggu PPK'."),
    ("4", "Verifikasi Status",
     "Setelah konfirmasi, refresh halaman. Status Tim Kerja tersebut berubah menjadi badge amber 'Menunggu PPK'."),
    ("5", "Ulangi untuk Tim Lain",
     "Lakukan hal yang sama untuk setiap Tim Kerja lain yang RA-nya berstatus 'Menunggu Review'."),
]
for i, (num, ttl, dsc) in enumerate(setujui_steps):
    step_row(slide, num, ttl, dsc, 6.7, 1.9 + i * 1.02, w=6.25, h=0.95, accent=HIJAU_MID)

note(slide, "ℹ  Setelah Kabag Umum menyetujui, dokumen diteruskan ke PPK. Kabag tidak perlu melakukan tindakan lanjutan.",
     6.7, 7.05, 6.25, 0.45, bg=RGBColor(0xE3, 0xF2, 0xFD), bd=BIRU_MUDA, fc=BIRU_TUA)

# ── Slide 8: Tolak RA ────────────────────────────────────────────────────────
slide = blank_slide(prs2)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "07  |  Menolak Rencana Aksi & Memberi Catatan Revisi",
       "Prosedur penolakan RA dengan catatan yang konstruktif")
footer(slide, 8, T2, "SIPITUNG — Juknis Kabag Umum | LLDIKTI Wilayah III")

ph(slide, "Tampilan Dialog Tolak & Catatan", 0.4, 1.35, 4.8, 5.75)

tb(slide, "Langkah Penolakan Rencana Aksi",
   5.55, 1.38, 7.4, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.55, 1.38, 0.06, 0.38, fill=AKSEN)

tolak_steps = [
    ("1", "Identifikasi Alasan Penolakan",
     "Sebelum menolak, pastikan Anda sudah mengetahui poin spesifik yang perlu diperbaiki oleh Tim Kerja."),
    ("2", "Klik Tombol 'Tolak'",
     "Klik tombol merah 'Tolak' pada kolom Aksi baris Tim Kerja yang bersangkutan."),
    ("3", "Dialog Catatan Revisi Terbuka",
     "Sebuah dialog berisi kolom teks 'Catatan Revisi' akan terbuka. Kolom ini WAJIB diisi sebelum dapat mengkonfirmasi penolakan."),
    ("4", "Tulis Catatan Revisi",
     "Tulis catatan revisi yang jelas, spesifik, dan konstruktif (contoh: 'Target TW II IKU 2.1 terlalu rendah dibanding target tahunan. Mohon disesuaikan.')."),
    ("5", "Klik Konfirmasi Tolak",
     "Klik tombol 'Konfirmasi Tolak'. Status dokumen berubah menjadi 'Ditolak' dan catatan terkirim ke Tim Kerja."),
    ("6", "Verifikasi",
     "Badge status Tim Kerja berubah menjadi merah 'Ditolak'. Tim Kerja dapat melihat catatan Anda dan melakukan perbaikan."),
]
for i, (num, ttl, dsc) in enumerate(tolak_steps):
    step_row(slide, num, ttl, dsc, 5.55, 1.9 + i * 0.85, w=7.4, h=0.78, accent=MERAH)

note(slide, "⚠  Penting: Catatan revisi harus SPESIFIK dan JELAS. Catatan yang terlalu umum (contoh: 'Perbaiki targetnya') akan menyulitkan Tim Kerja dalam melakukan revisi yang tepat.",
     5.55, 7.0, 7.4, 0.5)

# ── Slide 9: Review Laporan Pengukuran ───────────────────────────────────────
slide = blank_slide(prs2)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "08  |  Review Laporan Pengukuran Kinerja",
       "Cara mereview laporan triwulan yang disubmit Tim Kerja")
footer(slide, 9, T2, "SIPITUNG — Juknis Kabag Umum | LLDIKTI Wilayah III")

ph(slide, "Tampilan Halaman Laporan Pengukuran", 0.4, 1.35, 4.8, 5.75)

tb(slide, "Cara Mereview Laporan Pengukuran",
   5.55, 1.38, 7.4, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.55, 1.38, 0.06, 0.38, fill=AKSEN)

review_laporan = [
    ("1", "Buka Menu Pengukuran Kinerja",
     "Klik 'Pengukuran Kinerja' di sidebar. Halaman menampilkan kartu laporan per Tim Kerja di bagian atas, diikuti matriks realisasi."),
    ("2", "Pilih Periode Triwulan",
     "Pilih periode yang ingin direview dari dropdown. Kartu laporan akan berubah sesuai periode yang dipilih."),
    ("3", "Identifikasi Laporan Menunggu Review",
     "Cari kartu dengan badge biru 'Menunggu'. Kartu ini berisi tombol 'Setujui' dan 'Kembalikan'."),
    ("4", "Review Detail Realisasi",
     "Periksa matriks di bawah kartu laporan. Verifikasi nilai realisasi, progress kegiatan, kendala, dan strategi tindak lanjut setiap IKU."),
    ("5", "Cek Kelengkapan Semua IKU",
     "Pastikan semua IKU Tim Kerja sudah diisi (tidak ada baris kosong). Perhatikan apakah ada IKU kolaborasi yang belum terisi oleh pihak yang harusnya mengisi."),
    ("6", "Ambil Keputusan",
     "Jika data lengkap dan valid → Setujui. Jika ada yang perlu diperbaiki → Kembalikan dengan catatan rekomendasi."),
]
for i, (num, ttl, dsc) in enumerate(review_laporan):
    step_row(slide, num, ttl, dsc, 5.55, 1.9 + i * 0.85, w=7.4, h=0.78, accent=HIJAU_KAB2)

note(slide, "ℹ  Matriks realisasi menampilkan: Sasaran, IKU, PIC, Satuan, Target PK, Target TW, Realisasi, Progress, Kendala, Strategi, dan 'Diisi Oleh'.",
     5.55, 7.02, 7.4, 0.48, bg=RGBColor(0xE3, 0xF2, 0xFD), bd=BIRU_MUDA, fc=BIRU_TUA)

# ── Slide 10: Setujui/Tolak Laporan ──────────────────────────────────────────
slide = blank_slide(prs2)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "09  |  Setujui atau Tolak Laporan Pengukuran",
       "Langkah persetujuan dan prosedur penolakan laporan triwulan")
footer(slide, 10, T2, "SIPITUNG — Juknis Kabag Umum | LLDIKTI Wilayah III")

# Kiri — Setujui
rect(slide, 0.4, 1.35, 5.9, 5.75, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 0.4, 1.35, 5.9, 0.45, fill=HIJAU_MID)
tb(slide, "Menyetujui Laporan Pengukuran",
   0.55, 1.37, 5.65, 0.4, sz=Pt(12), bold=True, color=PUTIH)

setujui_laporan = [
    ("1", "Klik Tombol 'Setujui' pada Kartu Laporan",
     "Pada kartu laporan Tim Kerja yang sudah Anda review, klik tombol hijau 'Setujui'."),
    ("2", "Dialog Konfirmasi Muncul",
     "Sebuah dialog konfirmasi akan muncul. Anda dapat mengisi catatan apresiasi (opsional) sebelum mengkonfirmasi."),
    ("3", "Klik Konfirmasi",
     "Klik 'Konfirmasi'. Status laporan berubah menjadi 'Disetujui' (badge hijau). Tim Kerja mendapat notifikasi otomatis."),
    ("4", "Verifikasi",
     "Kartu laporan Tim Kerja tersebut sekarang menampilkan badge hijau 'Disetujui' beserta tanggal persetujuan."),
    ("5", "Ulangi untuk Tim Lain",
     "Lakukan hal yang sama untuk setiap laporan Tim Kerja lain yang perlu disetujui pada periode ini."),
]
for i, (num, ttl, dsc) in enumerate(setujui_laporan):
    step_row(slide, num, ttl, dsc, 0.4, 1.9 + i * 0.98, w=5.9, h=0.9, accent=HIJAU_MID)

# Kanan — Tolak/Kembalikan
rect(slide, 6.7, 1.35, 6.25, 5.75, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 6.7, 1.35, 6.25, 0.45, fill=MERAH)
tb(slide, "Menolak / Mengembalikan Laporan",
   6.85, 1.37, 6.0, 0.4, sz=Pt(12), bold=True, color=PUTIH)

tolak_laporan = [
    ("1", "Klik Tombol 'Kembalikan'",
     "Pada kartu laporan Tim Kerja, klik tombol merah 'Kembalikan'. Dialog catatan rekomendasi akan terbuka."),
    ("2", "Isi Catatan Rekomendasi (WAJIB)",
     "Tulis rekomendasi perbaikan secara spesifik (contoh: 'Realisasi IKU 3.1 perlu dilengkapi dengan keterangan progress kegiatan')."),
    ("3", "Klik Konfirmasi Kembalikan",
     "Klik 'Konfirmasi'. Status laporan berubah menjadi 'Dikembalikan'. Tim Kerja dapat melihat catatan dan melakukan perbaikan."),
    ("4", "Tim Kerja Melakukan Revisi",
     "Tim Kerja akan memperbaiki data realisasi dan submit ulang. Laporan yang disubmit ulang akan kembali berstatus 'Menunggu Review'."),
    ("5", "Review Ulang Laporan",
     "Setelah Tim Kerja submit ulang, ulangi proses review dari langkah awal. Tidak ada batasan jumlah siklus revisi."),
]
for i, (num, ttl, dsc) in enumerate(tolak_laporan):
    step_row(slide, num, ttl, dsc, 6.7, 1.9 + i * 0.98, w=6.25, h=0.9, accent=MERAH)

note(slide, "⚠  Pastikan catatan rekomendasi ditulis dengan bahasa yang konstruktif dan mudah dipahami. Ini akan membantu Tim Kerja melakukan perbaikan dengan cepat dan tepat.",
     0.4, 7.05, 12.5, 0.45)

# ── Slide 11: Export PDF ─────────────────────────────────────────────────────
slide = blank_slide(prs2)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "10  |  Export Laporan Pengukuran ke PDF",
       "Mengekspor laporan triwulan dalam format PDF untuk dokumentasi resmi")
footer(slide, 11, T2, "SIPITUNG — Juknis Kabag Umum | LLDIKTI Wilayah III")

ph(slide, "Tampilan Halaman Export PDF", 0.4, 1.35, 4.8, 5.75)

tb(slide, "Cara Export Laporan PDF",
   5.55, 1.38, 7.4, 0.38, sz=Pt(13), bold=True, color=BIRU_TUA)
rect(slide, 5.55, 1.38, 0.06, 0.38, fill=AKSEN)

export_steps = [
    ("1", "Buka Halaman Pengukuran Kinerja",
     "Klik 'Pengukuran Kinerja' di sidebar. Pastikan Anda sudah memilih periode triwulan yang ingin diekspor."),
    ("2", "Klik Tombol Export PDF",
     "Klik tombol 'Export PDF [TW]' yang tersedia di bagian atas halaman atau di area kartu laporan. Pilih triwulan yang sesuai."),
    ("3", "Halaman Preview Muncul",
     "Browser akan membuka halaman preview laporan dalam format A3 landscape. Periksa kelengkapan data sebelum mencetak."),
    ("4", "Cetak / Simpan sebagai PDF",
     "Gunakan fungsi cetak browser (Ctrl+P / Cmd+P). Pilih 'Save as PDF' sebagai printer untuk menyimpan, atau pilih printer fisik untuk mencetak."),
    ("5", "Verifikasi Isi Laporan",
     "Laporan PDF berisi: judul + periode, matriks realisasi (Sasaran, IKU, PIC, Target PK, Target TW, Realisasi), status laporan per Tim Kerja, dan area tanda tangan."),
]
for i, (num, ttl, dsc) in enumerate(export_steps):
    step_row(slide, num, ttl, dsc, 5.55, 1.9 + i * 0.95, w=7.4, h=0.88, accent=BIRU_MID)

# Isi laporan PDF
rect(slide, 5.55, 6.55, 7.4, 0.95, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 5.55, 6.55, 7.4, 0.08, fill=BIRU_MID)
tb(slide, "Konten dalam laporan PDF: Judul & Periode  |  Matriks Realisasi semua IKU  |  Status per Tim Kerja  |  Tanda Tangan",
   5.7, 6.68, 7.1, 0.75, sz=Pt(9.5), color=ABU_TUA, wrap=True)

# ── Slide 12: Checklist & Alur Lengkap ───────────────────────────────────────
slide = blank_slide(prs2)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "11  |  Alur Lengkap & Checklist Harian Kabag Umum",
       "Ringkasan alur kerja dan panduan checklist untuk memastikan tidak ada yang terlewat")
footer(slide, 12, T2, "SIPITUNG — Juknis Kabag Umum | LLDIKTI Wilayah III")

# Alur visual
rect(slide, 0.4, 1.35, 12.5, 1.55, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 0.4, 1.35, 12.5, 0.42, fill=BIRU_TUA)
tb(slide, "Alur Kerja Kabag Umum di SIPITUNG",
   0.55, 1.37, 12.2, 0.38, sz=Pt(11), bold=True, color=PUTIH)

alur2 = [
    ("Login &\nBuka Dashboard",          HIJAU_KAB2),
    ("Cek RA\nMenunggu Review",          BIRU_MID),
    ("Setujui /\nTolak RA",              HIJAU_MID),
    ("Cek Laporan\nPengukuran",          BIRU_TUA),
    ("Setujui /\nKembalikan",            HIJAU_KAB),
    ("Export\nPDF (opsional)",           BIRU_MUDA),
]
for i, (lbl, col) in enumerate(alur2):
    ax = 0.6 + i * 2.05
    rect(slide, ax, 1.85, 1.85, 1.0, fill=col)
    tb(slide, lbl, ax + 0.07, 1.95, 1.72, 0.82,
       sz=Pt(9.5), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
    if i < len(alur2)-1:
        tb(slide, "▶", ax + 1.87, 2.16, 0.17, 0.42,
           sz=Pt(11), color=BIRU_MUDA, align=PP_ALIGN.CENTER)

# Checklist dua kolom
rect(slide, 0.4, 3.05, 5.9, 4.0, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 0.4, 3.05, 5.9, 0.42, fill=HIJAU_KAB2)
tb(slide, "Checklist Review Rencana Aksi (RA)",
   0.55, 3.07, 5.65, 0.38, sz=Pt(12), bold=True, color=PUTIH)
cl_ra = [
    "Semua IKU sudah memiliki target TW I, TW II, TW III, TW IV",
    "Target triwulan bersifat kumulatif (TW IV = target tahunan)",
    "Tidak ada nilai target yang kosong atau nol tanpa keterangan",
    "IKU kolaborasi sudah di-assign ke PIC primer dan co-PIC yang benar",
    "Satuan IKU konsisten dengan dokumen PK yang telah disepakati",
    "Rencana kegiatan per IKU sudah terisi (minimal satu kegiatan per TW)",
]
for i, c in enumerate(cl_ra):
    ty = 3.57 + i * 0.57
    tb(slide, "☐", 0.55, ty + 0.03, 0.35, 0.44,
       sz=Pt(14), bold=True, color=HIJAU_MID, align=PP_ALIGN.CENTER)
    tb(slide, c, 0.95, ty + 0.06, 5.25, 0.46, sz=Pt(10), color=ABU_TUA, wrap=True)

rect(slide, 6.7, 3.05, 6.25, 4.0, fill=PUTIH,
     line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
rect(slide, 6.7, 3.05, 6.25, 0.42, fill=HIJAU_TUA)
tb(slide, "Checklist Review Laporan Pengukuran",
   6.85, 3.07, 6.0, 0.38, sz=Pt(12), bold=True, color=PUTIH)
cl_peng = [
    "Semua IKU Tim Kerja sudah diisi nilai realisasi (tidak ada yang kosong)",
    "Nilai realisasi wajar dan konsisten dengan periode sebelumnya",
    "Kolom 'Progress Kegiatan' terisi dengan deskripsi yang jelas",
    "Kolom 'Kendala' diisi jika ada hambatan yang signifikan",
    "Kolom 'Strategi Tindak Lanjut' diisi sesuai kendala yang ada",
    "Untuk IKU kolaborasi: terisi oleh pihak yang tepat (cek kolom 'Diisi Oleh')",
]
for i, c in enumerate(cl_peng):
    ty = 3.57 + i * 0.57
    tb(slide, "☐", 6.85, ty + 0.03, 0.35, 0.44,
       sz=Pt(14), bold=True, color=HIJAU_MID, align=PP_ALIGN.CENTER)
    tb(slide, c, 7.25, ty + 0.06, 5.6, 0.46, sz=Pt(10), color=ABU_TUA, wrap=True)

# ── Slide 13: FAQ ────────────────────────────────────────────────────────────
slide = blank_slide(prs2)
rect(slide, 0, 0, 13.33, 7.5, fill=ABU_MUDA)
header(slide, "12  |  FAQ — Kabag Umum", "Pertanyaan yang Sering Diajukan")
footer(slide, 13, T2, "SIPITUNG — Juknis Kabag Umum | LLDIKTI Wilayah III")

faqs2 = [
    ("Tombol Setujui/Tolak tidak muncul di baris Tim Kerja tertentu.",
     "Tombol hanya muncul jika status dokumen adalah 'submitted'. Jika status lain (draft, kabag_approved, dll.), kolom Aksi menampilkan tanda '—'."),
    ("Saya sudah menyetujui RA, tapi tim kerja bilang belum berubah statusnya.",
     "Minta Tim Kerja untuk me-refresh halaman (F5 atau Ctrl+R). Status diperbarui secara real-time setelah persetujuan."),
    ("Apakah saya bisa membatalkan persetujuan yang sudah saya berikan?",
     "Tidak. Setelah Kabag Umum menyetujui, dokumen diteruskan ke PPK. Hanya SuperAdmin yang dapat membuka kembali dokumen yang sudah ppk_approved."),
    ("Saya tidak bisa melihat dokumen salah satu Tim Kerja.",
     "Pastikan filter tahun anggaran sudah benar. Jika masih tidak tampil, kemungkinan Tim Kerja tersebut belum membuat dokumen (status masih draft). Hubungi Tim Kerja terkait."),
    ("Bagaimana cara mengetahui laporan mana saja yang masih menunggu review?",
     "Lihat Dashboard Kabag Umum — angka 'Menunggu Review' menunjukkan jumlah dokumen yang perlu ditindaklanjuti. Klik untuk langsung menuju halaman tersebut."),
    ("Apakah saya bisa export laporan sebelum semua Tim Kerja submit?",
     "Ya, export bisa dilakukan kapan saja. Namun, laporan yang belum disubmit Tim Kerja tidak akan muncul dalam data realisasi yang diekspor."),
]
for i, (q, a) in enumerate(faqs2):
    col = i % 2; row = i // 2
    lx = 0.4 + col * 6.5; ty = 1.38 + row * 1.92
    rect(slide, lx, ty, 6.1, 1.85, fill=PUTIH,
         line=RGBColor(0xCC, 0xD9, 0xEA), lw=Pt(1))
    rect(slide, lx, ty, 6.1, 0.08, fill=HIJAU_KAB2)
    rect(slide, lx, ty + 0.08, 0.52, 1.77, fill=RGBColor(0xD5, 0xF5, 0xE3))
    tb(slide, "Q", lx + 0.1, ty + 0.66, 0.34, 0.44,
       sz=Pt(18), bold=True, color=HIJAU_KAB, align=PP_ALIGN.CENTER)
    tb(slide, q, lx + 0.62, ty + 0.1, 5.35, 0.48,
       sz=Pt(10), bold=True, color=BIRU_TUA, wrap=True)
    rect(slide, lx + 0.62, ty + 0.6, 5.35, 0.02, fill=RGBColor(0xCC, 0xD9, 0xEA))
    tb(slide, a, lx + 0.62, ty + 0.68, 5.35, 1.12,
       sz=Pt(9.5), color=ABU_TUA, wrap=True)

# ── Slide 14: Penutup ────────────────────────────────────────────────────────
slide = blank_slide(prs2)
rect(slide, 0, 0, 13.33, 7.5, fill=BIRU_TUA)
rect(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0x0F, 0x20, 0x44))
rect(slide, 0, 3.5, 13.33, 0.08, fill=AKSEN)

tb(slide, "Terima Kasih", 1.0, 1.5, 11.33, 1.2,
   sz=Pt(44), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)
tb(slide, "Juknis SIPITUNG — Kabag Umum | LLDIKTI Wilayah III",
   1.0, 2.8, 11.33, 0.6, sz=Pt(16), color=AKSEN, align=PP_ALIGN.CENTER)

tb(slide, "Untuk pertanyaan dan bantuan teknis, hubungi Admin Sistem SIPITUNG",
   1.0, 3.8, 11.33, 0.5, sz=Pt(13), color=PUTIH, align=PP_ALIGN.CENTER)
tb(slide, "Versi 1.0  |  April 2026", 1.0, 4.4, 11.33, 0.45,
   sz=Pt(11), color=RGBColor(0x88, 0xA8, 0xD0), align=PP_ALIGN.CENTER)

modul_tags = [("Perencanaan", HIJAU_KAB2), ("Rencana Aksi", BIRU_MID), ("Pengukuran Kinerja", HIJAU_MID)]
for i, (tag, col) in enumerate(modul_tags):
    tx = 3.3 + i * 2.3
    rect(slide, tx, 5.2, 2.1, 0.55, fill=col)
    tb(slide, tag, tx + 0.1, 5.26, 1.92, 0.42,
       sz=Pt(11), bold=True, color=PUTIH, align=PP_ALIGN.CENTER)

# Simpan file 2
out2 = "Juknis_SIPITUNG_Kabag_Umum.pptx"
prs2.save(out2)
print(f"✓ File 2 selesai: {out2}  ({len(prs2.slides)} slides)")
print("\nSemua file berhasil dibuat.")
